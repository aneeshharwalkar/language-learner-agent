"""
Lightning Talk deck — Language Learner Agent
8-min talk · 11 slides · presenter-driven · Tract dark theme

Design rules (from Preston):
- Sparse slides. 1 idea per slide. Big type. Speaker fills in the rest.
- No meta/rubric language on slide. No em dashes. Plain English glosses.
- Match research deck's Tract dark palette.
- Aunt hook is verbal (in speaker notes), not on the slide.
- Speaker notes carry the full script + speaker label + timing.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

# ---- Tokens ----
BG       = "0B0D10"
SURFACE  = "13161B"
CARD     = "1A1E26"
RULE     = "282C35"
TEXT     = "F2F3F5"
MUTED    = "9AA1AB"
DIM      = "656B75"

BLUE  = "7B9FE0"   # ACT I · EVALUATE
PINK  = "E07BAA"   # ACT II · ANALYZE  /  L3
GREEN = "7BE0B4"   # ACT III · IMPROVE  / close
AMBER = "E0C27B"   # L2
RED   = "E07B7B"   # L1

PHASE_COLORS = {"I": BLUE, "II": PINK, "III": GREEN}

FONT_HEAD = "Helvetica Neue"
FONT_MONO = "Menlo"

def emu(v): return int(v * 914400)

# ---- Helpers ----
def set_bg(slide):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor.from_string(BG)

def add_text(slide, text, l, t, w, h, *, size=14, color=TEXT, bold=False, italic=False,
             font=FONT_HEAD, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(emu(l), emu(t), emu(w), emu(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = RGBColor.from_string(color)
    return tb

def add_rect(slide, l, t, w, h, fill_hex, *, line=False):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, emu(l), emu(t), emu(w), emu(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = RGBColor.from_string(fill_hex)
    if not line:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def add_eyebrow(slide, label, color_hex):
    add_text(slide, label, 0.60, 0.45, 7.50, 0.30,
             size=11, color=color_hex, bold=True, font=FONT_HEAD)
    # rule from end of label to right edge
    add_rect(slide, 5.60, 0.56, 7.20, 0.01, RULE)

def add_pagenum(slide, n):
    add_text(slide, str(n), 12.30, 7.12, 0.80, 0.30,
             size=9, color=DIM, align=PP_ALIGN.RIGHT)

def add_notes(slide, text):
    nf = slide.notes_slide.notes_text_frame
    nf.text = text

def add_blank_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    return s

# ---- Build ----
prs = Presentation()
prs.slide_width = emu(13.33)
prs.slide_height = emu(7.50)


# =========================================================
# SLIDE 1 — Title (P, 0:15)
# =========================================================
s = add_blank_slide(prs)

# Eyebrow (small, dim — sets tone without competing with title)
add_text(s, "APPLIED HUMAN FACTORS · SPRING 2026", 0.60, 0.45, 8.00, 0.30,
         size=10, color=MUTED, bold=True)
add_rect(s, 7.40, 0.56, 5.40, 0.01, RULE)

# Title (centered, big — sized to fit on one line)
add_text(s, "Bridging the Learning-Transfer Gap", 0.60, 2.55, 12.13, 1.10,
         size=44, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
add_text(s, "Deploying conversational agents in Duolingo.", 0.60, 3.75, 12.13, 0.55,
         size=20, color=MUTED, align=PP_ALIGN.CENTER)

# Bottom team + arc
add_text(s, "PRESTON CUSICK · ANEESH HARWALKAR", 0.60, 5.90, 12.13, 0.40,
         size=12, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, "Evaluate  →  Analyze  →  Improve", 0.60, 6.35, 12.13, 0.40,
         size=14, color=DIM, italic=True, align=PP_ALIGN.CENTER)

add_pagenum(s, 1)
add_notes(s, """SPEAKER: Preston · 0:15

[AUNT HOOK — told to the room before clicking forward]

"My aunt kept a Duolingo streak for three years. Italian, every day. Last summer she finally went to Italy. She stood in front of her own family. And she couldn't speak to them. Three years of correct answers. No conversation. That gap, right there, is our project."

[Click to slide 2]
""")


# =========================================================
# SLIDE 2 — The System (P, 0:50)
# =========================================================
s = add_blank_slide(prs)
add_eyebrow(s, "ACT I · EVALUATE", BLUE)

add_text(s, "500M+ monthly active users.", 0.60, 0.95, 12.13, 1.20,
         size=44, color=TEXT, bold=True)
add_text(s, "One of the largest human-machine learning systems on Earth.", 0.60, 2.10, 12.13, 0.50,
         size=18, color=MUTED)

# Three stressor columns
cols = [
    ("INTERNAL",       "Cognitive load. The illusion that recognizing a word means you know it.",      BLUE),
    ("EXTERNAL",       "Streaks. XP. Gamification that rewards showing up, not speaking.",             PINK),
    ("ORGANIZATIONAL", "Daily-active-user incentives. Innovator's Dilemma locks in the wrong metric.", AMBER),
]
col_w, gap = 4.00, 0.20
total_w = col_w * 3 + gap * 2
left0 = (13.33 - total_w) / 2  # center horizontally
top, h = 3.50, 2.60
for i, (label, body, color) in enumerate(cols):
    L = left0 + i * (col_w + gap)
    add_rect(s, L, top, col_w, h, CARD)
    # accent bar
    add_rect(s, L, top, 0.06, h, color)
    add_text(s, label, L + 0.30, top + 0.30, col_w - 0.50, 0.40,
             size=11, color=color, bold=True)
    add_text(s, body, L + 0.30, top + 0.85, col_w - 0.50, h - 1.00,
             size=14, color=TEXT)

add_pagenum(s, 2)
add_notes(s, """SPEAKER: Preston · 0:50

"Duolingo is a human-machine system with over 500 million monthly active users. From Week 1 of class, every system sits under three kinds of stressors.

Internal: cognitive load, and the illusion that seeing a word again means you know it.

External: streaks, XP, the stuff the app makes you feel bad about when you miss a day.

Organizational: Duolingo's business rewards daily active users. Christensen's Innovator's Dilemma — they can't rebuild around real conversation without breaking the metric paying the bills."

[~50 seconds. Hand off to Aneesh.]
""")


# =========================================================
# SLIDE 3 — The Ceiling (A, 0:50)
# =========================================================
s = add_blank_slide(prs)
add_eyebrow(s, "ACT I · EVALUATE", BLUE)

add_text(s, "The measurable ceiling.", 0.60, 0.95, 12.13, 1.10,
         size=44, color=TEXT, bold=True)
add_text(s, "Jiang, Rollinson, Plonsky et al. (2021). ACTFL proficiency scale.",
         0.60, 2.05, 12.13, 0.45,
         size=14, color=MUTED, italic=True)

# Three big stat cards
stats = [
    ("READING",  "Intermediate",  "Low",      BLUE),
    ("LISTENING","Novice",        "High",     AMBER),
    ("SPEAKING", "not assessed", "",          RED),
]
col_w, gap = 3.95, 0.25
total_w = col_w * 3 + gap * 2
left0 = (13.33 - total_w) / 2
top, h = 3.10, 2.80
for i, (label, big, small, color) in enumerate(stats):
    L = left0 + i * (col_w + gap)
    add_rect(s, L, top, col_w, h, CARD)
    add_rect(s, L, top, 0.06, h, color)
    add_text(s, label, L + 0.30, top + 0.35, col_w - 0.50, 0.40,
             size=11, color=color, bold=True)
    add_text(s, big, L + 0.30, top + 1.00, col_w - 0.50, 0.85,
             size=30, color=TEXT, bold=True)
    if small:
        add_text(s, small, L + 0.30, top + 1.85, col_w - 0.50, 0.55,
                 size=24, color=MUTED)

# Pull quote
add_text(s, '"No other skills were assessed."', 0.60, 6.20, 12.13, 0.50,
         size=18, color=MUTED, italic=True, align=PP_ALIGN.CENTER)

add_pagenum(s, 3)
add_notes(s, """SPEAKER: Aneesh · 0:50

"Jiang, Rollinson, and Plonsky put Duolingo learners through the ACTFL proficiency scale. Reading topped out at Intermediate Low. Listening at Novice High. Speaking — speaking was never tested.

In the peer-reviewed outcome research we surveyed, the paper itself says: 'no other skills were assessed.'

So the ceiling is real, it's measurable, and the productive skills, the ones that let you actually have a conversation, are invisible to the system that claims to be teaching you."

[~50 seconds.]
""")


# =========================================================
# SLIDE 4 — Task & User Demands (A, 0:45)
# =========================================================
s = add_blank_slide(prs)
add_eyebrow(s, "ACT II · ANALYZE", PINK)

add_text(s, "Four lenses on the failure.", 0.60, 0.95, 12.13, 1.10,
         size=44, color=TEXT, bold=True)
add_text(s, "The class demand framework: where the system loads and unloads the learner.",
         0.60, 2.05, 12.13, 0.45, size=14, color=MUTED, italic=True)

quads = [
    ("ATTENTION",
     "Learner gaze fixates on XP and streaks. Inattentional blindness to the linguistic content.",
     BLUE),
    ("MEMORY",
     "Phonological Loop is exercised. Episodic Buffer is never engaged. Words don't bind to context.",
     PINK),
    ("WORKLOAD",
     "Over-loaded on gamification. Under-loaded on speech production.",
     AMBER),
    ("COGNITION",
     "No real-world knowledge built. No regional speech patterns. No conversation skills.",
     GREEN),
]
qw, gap_x, gap_y = 6.00, 0.20, 0.20
qh = 1.85
left0 = (13.33 - (qw * 2 + gap_x)) / 2
top0 = 2.85
for i, (label, body, color) in enumerate(quads):
    row, col = i // 2, i % 2
    L = left0 + col * (qw + gap_x)
    T = top0 + row * (qh + gap_y)
    add_rect(s, L, T, qw, qh, CARD)
    add_rect(s, L, T, 0.06, qh, color)
    add_text(s, label, L + 0.30, T + 0.25, qw - 0.50, 0.40,
             size=11, color=color, bold=True)
    add_text(s, body, L + 0.30, T + 0.75, qw - 0.50, qh - 0.85,
             size=14, color=TEXT)

add_pagenum(s, 4)
add_notes(s, """SPEAKER: Aneesh · 0:45

"We ran the four HF demand lenses from class.

Attention: the learner's gaze and mental budget go to XP bars and streaks, not to language.

Memory: Duolingo exercises Baddeley's Phonological Loop — verbal rehearsal — but never the Episodic Buffer, the part of working memory that binds words to context.

Workload: over-loaded on gamification. Under-loaded on speech production.

Cognition: the real-world knowledge you need for actual conversation is never built.

The recognition task succeeds. The production task fails."

[~45 seconds.]
""")


# =========================================================
# SLIDE 5 — Methodology (P, 0:40)
# =========================================================
s = add_blank_slide(prs)
add_eyebrow(s, "ACT II · ANALYZE", PINK)

add_text(s, "How we got here.", 0.60, 0.95, 12.13, 1.10,
         size=44, color=TEXT, bold=True)
add_text(s, "Two-week timeline. No primary data. We are upfront about that.",
         0.60, 2.05, 12.13, 0.45, size=14, color=MUTED, italic=True)

# 3 method cards
methods = [
    ("EVIDENCE BASE", "15 peer-reviewed sources.", "Spans HF, SLA, working memory, and dialogue-based learning.", BLUE),
    ("EFFECT SIZES",  "Two meta-analyses carry the claims.", "Bibauw 2022 · d = 0.58.   Mayer 2021 · d = 1.3.",      PINK),
    ("PROPOSED RCT",  "24-week, dialect-matched.",            "Extends Wang 2024. Speaking performance is the outcome.", GREEN),
]
col_w, gap = 4.00, 0.20
total_w = col_w * 3 + gap * 2
left0 = (13.33 - total_w) / 2
top, h = 3.00, 3.20
for i, (label, headline, body, color) in enumerate(methods):
    L = left0 + i * (col_w + gap)
    add_rect(s, L, top, col_w, h, CARD)
    add_rect(s, L, top, 0.06, h, color)
    add_text(s, label, L + 0.30, top + 0.30, col_w - 0.50, 0.40,
             size=11, color=color, bold=True)
    add_text(s, headline, L + 0.30, top + 0.85, col_w - 0.50, 1.00,
             size=20, color=TEXT, bold=True)
    add_text(s, body, L + 0.30, top + 1.95, col_w - 0.50, h - 2.10,
             size=13, color=MUTED)

add_text(s, "Cohen's d is a standardized effect size. d = 0.5 is medium, d ≥ 0.8 is large.",
         0.60, 6.50, 12.13, 0.35, size=11, color=DIM, italic=True, align=PP_ALIGN.CENTER)

add_pagenum(s, 5)
add_notes(s, """SPEAKER: Preston · 0:40

"Two-week timeline, so we are upfront: the biggest threat to validity is no primary data.

Our approach is research-and-build. Fifteen peer-reviewed sources. Two meta-analyses carrying the effect-size claims: Bibauw at d = 0.58 for dialogue-based learning, Mayer at d = 1.3 for personalized agents.

For causal evidence, we propose a 24-week randomized controlled trial that extends Wang's 2024 design, with dialect-matched conditions and speaking performance as the dependent measure."

[~40 seconds. Hand off to Aneesh.]
""")


# =========================================================
# SLIDE 6 — HFACS · The Method (A, 0:35)
# =========================================================
s = add_blank_slide(prs)
add_eyebrow(s, "ACT II · ANALYZE", PINK)

add_text(s, "Reason's Swiss Cheese.", 0.60, 0.95, 12.13, 1.10,
         size=44, color=TEXT, bold=True)
add_text(s, "Four layers of defense. Don't blame the operator. Work backward through all of them.",
         0.60, 2.05, 12.13, 0.45, size=14, color=MUTED, italic=True)

# Reason's Swiss Cheese — 4 layers with small round holes.
# One hole column is aligned vertically across all 4 layers; a red
# "failure trajectory" arrow pierces through it. Iconic to the model.
ALIGNED_X = 8.40   # the x offset (from cheese_L) where holes line up
layers = [
    # (label, accent, list of hole x-offsets for that layer)
    ("L4 · ORGANIZATIONAL INFLUENCES", BLUE,  [4.20, 6.10, ALIGNED_X]),
    ("L3 · UNSAFE SUPERVISION",        PINK,  [5.00, ALIGNED_X, 9.60]),
    ("L2 · PRECONDITIONS",             AMBER, [3.40, ALIGNED_X, 7.20]),
    ("L1 · UNSAFE ACTS",               RED,   [4.60, 6.40, ALIGNED_X]),
]
cheese_L = 1.20
cheese_W = 11.00
layer_h = 0.70
gap = 0.18
top0 = 2.85
hole_d = 0.34   # circle diameter

for i, (label, color, hole_xs) in enumerate(layers):
    T = top0 + i * (layer_h + gap)
    add_rect(s, cheese_L, T, cheese_W, layer_h, CARD)
    add_rect(s, cheese_L, T, 0.06, layer_h, color)
    add_text(s, label, cheese_L + 0.30, T + 0.15, 7.00, 0.40,
             size=12, color=color, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    # round holes (BG fill so they read as "punched through")
    for hx in hole_xs:
        hole = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                  emu(cheese_L + hx),
                                  emu(T + (layer_h - hole_d) / 2),
                                  emu(hole_d), emu(hole_d))
        hole.fill.solid()
        hole.fill.fore_color.rgb = RGBColor.from_string(BG)
        hole.line.fill.background()
        hole.shadow.inherit = False

# Failure trajectory: red vertical arrow piercing the aligned hole column
arrow_x_center = cheese_L + ALIGNED_X + hole_d / 2
arrow_w = 0.10
# vertical shaft — from above the top bar to just into the bottom bar
shaft_top = top0 - 0.45
shaft_bot = top0 + 4 * layer_h + 3 * gap - 0.20
shaft = add_rect(s, arrow_x_center - arrow_w / 2, shaft_top, arrow_w, shaft_bot - shaft_top, RED)

# Arrowhead (downward triangle) at the bottom
arrowhead = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                               emu(arrow_x_center - 0.22),
                               emu(shaft_bot),
                               emu(0.44), emu(0.30))
arrowhead.rotation = 180
arrowhead.fill.solid()
arrowhead.fill.fore_color.rgb = RGBColor.from_string(RED)
arrowhead.line.fill.background()
arrowhead.shadow.inherit = False

# "HAZARD" label above the arrow
add_text(s, "HAZARD", arrow_x_center - 1.00, shaft_top - 0.40, 2.00, 0.30,
         size=10, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
# "INCIDENT" label below the arrowhead
add_text(s, "INCIDENT", arrow_x_center - 1.00, shaft_bot + 0.35, 2.00, 0.30,
         size=10, color=RED, bold=True, align=PP_ALIGN.CENTER)

add_pagenum(s, 6)
add_notes(s, """SPEAKER: Aneesh · 0:35

"For the Human Factors analysis itself we used HFACS — the Human Factors Analysis and Classification System, built on Reason's Swiss Cheese model.

Four layers of defense. Level 4: Organizational Influences. Level 3: Unsafe Supervision. Level 2: Preconditions. Level 1: Unsafe Acts at the operator.

The course's core tenet: work backward through all four. When the holes line up, an incident gets through. Don't blame the person — look for where the system failed."

[~35 seconds. Set up the star slide.]
""")


# =========================================================
# SLIDE 7 — The Holes Line Up (A, 1:10) ★ STAR SLIDE
# =========================================================
s = add_blank_slide(prs)
add_eyebrow(s, "ACT II · ANALYZE", PINK)

add_text(s, "The holes line up.", 0.60, 0.95, 12.13, 1.10,
         size=44, color=TEXT, bold=True)
add_text(s, "Each level fails into the next. The incident is structurally predicted.",
         0.60, 2.05, 12.13, 0.45, size=14, color=MUTED, italic=True)

# 4 stacked rows · L1 → L4 top to bottom (matches research deck slide 25)
rows = [
    ("L1 · UNSAFE ACTS",               "The learner sees XP. Gets no real signal on speaking.",       RED),
    ("L2 · PRECONDITIONS",             "Episodic Buffer never engaged. Words don't bind to context.", AMBER),
    ("L3 · UNSAFE SUPERVISION",        "Birdbrain optimizes item-correctness, not conversation.",     PINK),
    ("L4 · ORGANIZATIONAL INFLUENCES", "DAU metric excludes speaking entirely.",                       BLUE),
]
row_L, row_W = 1.20, 11.00
row_h, gap = 0.85, 0.15
top0 = 2.80
for i, (label, body, color) in enumerate(rows):
    T = top0 + i * (row_h + gap)
    add_rect(s, row_L, T, row_W, row_h, CARD)
    add_rect(s, row_L, T, 0.06, row_h, color)
    add_text(s, label, row_L + 0.30, T + 0.15, 4.20, 0.55,
             size=12, color=color, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, body, row_L + 4.60, T + 0.15, row_W - 4.90, 0.55,
             size=14, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)

add_pagenum(s, 7)
add_notes(s, """SPEAKER: Aneesh · 1:10  ★ STAR SLIDE — earns the Systems Thinking points

"Walk it with me, top to bottom.

Level 1: the learner acts on the only signal available, which is XP. In Signal Detection terms, every streak is a False Alarm on the variable that matters.

Level 2: without a supervisor asking for speech, the environment never engages the Episodic Buffer, and words never bind to context. Essoe's working-memory study showed this cleanly: distinct contexts produced 92% retention versus 76% in uniform contexts.

Level 3: Birdbrain, Duolingo's own algorithm, optimizes item-correctness. By TADMUS principles, the supervisor's sensitivity, d-prime, on conversational failure is effectively zero.

Level 4: the organization measures Daily Active Users, which explicitly excludes speaking.

The holes line up. This isn't a user mistake. The incident is structurally predicted."

[~70 seconds. This is the rigor + systems-thinking slide.]
""")


# =========================================================
# SLIDE 8 — Name the Enemy (P, 0:40)
# =========================================================
s = add_blank_slide(prs)
add_eyebrow(s, "ACT II · ANALYZE", PINK)

add_text(s, "Name the enemy.", 0.60, 0.95, 12.13, 1.10,
         size=44, color=TEXT, bold=True)
add_text(s, "Rosling. Tufte. Musk. Make the failure concrete.",
         0.60, 2.05, 12.13, 0.45, size=14, color=MUTED, italic=True)

enemies = [
    ("01", "THE ISOLATED DRILL",            "Vocabulary taught as isolated items, never embedded in a real situation.", RED),
    ("02", "THE WRONG SIGNAL",              "Confident feedback about XP and streaks. The wrong variable, amplified.",  AMBER),
    ("03", "THE RECEPTIVE-ONLY EVIDENCE",   "Outcome research covers reading and listening. Speaking is unsupervised.", PINK),
    ("04", "THE ALIGNED HOLES",             "Fix one layer alone and the cheese still lines up. The redesign spans levels.", BLUE),
]
qw, gap_x, gap_y = 6.00, 0.20, 0.20
qh = 1.95
left0 = (13.33 - (qw * 2 + gap_x)) / 2
top0 = 2.80
for i, (num, title, body, color) in enumerate(enemies):
    row, col = i // 2, i % 2
    L = left0 + col * (qw + gap_x)
    T = top0 + row * (qh + gap_y)
    add_rect(s, L, T, qw, qh, CARD)
    add_rect(s, L, T, 0.06, qh, color)
    # number
    add_text(s, num, L + 0.30, T + 0.30, 0.80, 0.55,
             size=20, color=color, bold=True)
    # title
    add_text(s, title, L + 1.20, T + 0.32, qw - 1.40, 0.50,
             size=14, color=color, bold=True)
    # body
    add_text(s, body, L + 1.20, T + 0.90, qw - 1.40, qh - 1.00,
             size=13, color=TEXT)

add_pagenum(s, 8)
add_notes(s, """SPEAKER: Preston · 0:40

"The course says: name the enemy. Rosling, Tufte, Musk. So here they are, named.

One: the Isolated Drill. Vocabulary taught as isolated items, not embedded in any real situation.

Two: the Wrong Signal. Confident feedback about XP and streaks. The wrong variable, amplified.

Three: the Receptive-Only Evidence Base. Duolingo's own research measures reading and listening. Speaking is unsupervised by design.

Four: the Aligned Holes. Fix any single layer in isolation and the cheese still lines up. Our redesign has to span levels."

[~40 seconds. Hand back to Preston for redesign.]
""")


# =========================================================
# SLIDE 9 — The Redesign · Carlos & Elena (P, 1:10) ★ DESIGN SLIDE
# =========================================================
s = add_blank_slide(prs)
add_eyebrow(s, "ACT III · IMPROVE", GREEN)

add_text(s, "Carlos. Elena. Real dialect, real conversation.",
         0.60, 0.95, 12.13, 1.10, size=40, color=TEXT, bold=True)
add_text(s, "A dialect-specific conversational agent that scaffolds the interaction.",
         0.60, 2.05, 12.13, 0.45, size=14, color=MUTED, italic=True)

# Two phone screenshots side by side on left half
img_top = 2.65
img_h = 4.50
# Carlos
img_w_carlos = img_h * 0.541  # 786/1454 ≈ 0.541
s.shapes.add_picture("/tmp/deck_audit/carlos.png",
                     emu(0.80), emu(img_top),
                     width=emu(img_w_carlos), height=emu(img_h))
# Elena
img_w_elena = img_h * 0.563  # 802/1424 ≈ 0.563
s.shapes.add_picture("/tmp/deck_audit/elena.png",
                     emu(0.80 + img_w_carlos + 0.30), emu(img_top),
                     width=emu(img_w_elena), height=emu(img_h))

# Three pillars, right side
pillars = [
    ("DIALOGUE WORKS",  "Bibauw 2022. d = 0.58 across 17 studies.",                    BLUE),
    ("PERSONA MATTERS", "Mayer 2021. d = 1.3 across 10 of 10 experiments.",            PINK),
    ("DIALECT IS INPUT","Geeslin and Long 2014. Real regional speech patterns matter.", GREEN),
]
pl_L = 6.40
pl_W = 6.30
pl_H = 1.35
pl_gap = 0.15
pl_top0 = img_top
for i, (label, body, color) in enumerate(pillars):
    T = pl_top0 + i * (pl_H + pl_gap)
    add_rect(s, pl_L, T, pl_W, pl_H, CARD)
    add_rect(s, pl_L, T, 0.06, pl_H, color)
    add_text(s, label, pl_L + 0.30, T + 0.25, pl_W - 0.50, 0.40,
             size=12, color=color, bold=True)
    add_text(s, body, pl_L + 0.30, T + 0.72, pl_W - 0.50, pl_H - 0.85,
             size=15, color=TEXT)

add_text(s, "Carlos · Mexican Spanish.   Elena · Castilian Spanish.",
         0.60, 7.10, 6.00, 0.30, size=10, color=DIM, italic=True)

add_pagenum(s, 9)
add_notes(s, """SPEAKER: Preston · 1:10  ★ DESIGN / ARCHITECTURAL INSIGHT slide

"Here's the proof-of-concept. A dialect-specific conversational agent — Carlos, Mexican Spanish, and Elena, Castilian — that scaffolds the interaction rather than just responding.

It rests on three findings.

One: Bibauw's 2022 meta-analysis across seventeen studies. Dialogue-based learning produces a medium effect, d equals point-five-eight, on second-language proficiency.

Two: Mayer's Personalization Principle. d equals one-point-three across ten of ten experiments. Embodiment and conversational style aren't UX decoration — they're the mechanism.

Three: Geeslin and Long. Communicative competence requires acquiring real regional speech patterns. Classrooms don't provide that. Duolingo doesn't provide that. That's our wedge."

[~70 seconds. Hand off to Aneesh.]
""")


# =========================================================
# SLIDE 10 — System Map + Honest Gaps (A, 0:55)
# =========================================================
s = add_blank_slide(prs)
add_eyebrow(s, "ACT III · IMPROVE", GREEN)

add_text(s, "Before. After. What's still missing.",
         0.60, 0.95, 12.13, 1.10, size=40, color=TEXT, bold=True)

# Top half: two horizontal flow diagrams
flow_top = 2.40
flow_h = 1.90

# BEFORE
add_rect(s, 0.60, flow_top, 6.00, flow_h, CARD)
add_rect(s, 0.60, flow_top, 0.06, flow_h, RED)
add_text(s, "BEFORE", 0.90, flow_top + 0.20, 5.50, 0.35,
         size=11, color=RED, bold=True)
add_text(s, "Learner  →  Item  →  XP", 0.90, flow_top + 0.70, 5.50, 0.50,
         size=18, color=TEXT, bold=True)
add_text(s, "Loop closes on the wrong variable.",
         0.90, flow_top + 1.30, 5.50, 0.45,
         size=12, color=MUTED, italic=True)

# AFTER
add_rect(s, 6.80, flow_top, 6.00, flow_h, CARD)
add_rect(s, 6.80, flow_top, 0.06, flow_h, GREEN)
add_text(s, "AFTER", 7.10, flow_top + 0.20, 5.50, 0.35,
         size=11, color=GREEN, bold=True)
add_text(s, "Learner  →  Dialect Agent  →  Speaking",
         7.10, flow_top + 0.70, 5.50, 0.50, size=18, color=TEXT, bold=True)
add_text(s, "Feedback on production, not recognition.",
         7.10, flow_top + 1.30, 5.50, 0.45,
         size=12, color=MUTED, italic=True)

# Bottom: honest gaps
add_text(s, "Honest gaps.", 0.60, 4.70, 12.13, 0.45,
         size=18, color=AMBER, bold=True)

gaps = [
    ("No primary data.",                     "Two-week window."),
    ("No A/B vs. a generic agent.",          "Dialect-effect is theorized, not tested."),
    ("Wang 2024 moved affect, not skill.",   "Closest comparable RCT. Plausibility, not proof."),
]
gtop = 5.25
gh = 0.55
for i, (title, sub) in enumerate(gaps):
    T = gtop + i * gh
    add_text(s, "•", 0.60, T, 0.30, gh, size=14, color=AMBER, bold=True)
    add_text(s, title, 0.95, T, 5.50, gh, size=14, color=TEXT, bold=True)
    add_text(s, sub, 6.50, T, 6.50, gh, size=14, color=MUTED)

add_text(s, "NEXT  ·  Eye-tracking at UT's IX Lab to sharpen the Level 1 attention claim.",
         0.60, 7.05, 12.13, 0.35,
         size=12, color=GREEN, italic=True, align=PP_ALIGN.CENTER)

add_pagenum(s, 10)
add_notes(s, """SPEAKER: Aneesh · 0:55

"Here's the system map. Before, Duolingo's loop runs learner to item to XP, and the loop closes on the wrong variable. After, the agent sits in the middle, dialogue scaffolds context, and feedback is about production — speaking — not recognition.

Three honest gaps. No primary data. No A/B comparing dialect-specific to generic agents. And Wang 2024 — the closest adult-learner GenAI chatbot study — moved willingness to communicate over twelve weeks, but it did not move measured speaking performance.

A proof-of-concept is proof of plausibility, not proof of effect. Our next step is an eye-tracking study at UT's IX Lab to turn the Level 1 attention claim from conceptual to measurable."

[~55 seconds. Hand back to Preston for close.]
""")


# =========================================================
# SLIDE 11 — Close · The Tenet (P, 0:25)
# =========================================================
s = add_blank_slide(prs)
add_eyebrow(s, "PROJECT CLOSE", GREEN)

# Big single line, centered vertically
add_text(s, '"Error is a system property,', 0.60, 2.40, 12.13, 1.00,
         size=42, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
add_text(s, 'not a human shortcoming."', 0.60, 3.30, 12.13, 1.00,
         size=42, color=TEXT, bold=True, align=PP_ALIGN.CENTER)

add_text(s, "System failure, not user failure.",
         0.60, 4.65, 12.13, 0.50, size=18, color=MUTED, italic=True, align=PP_ALIGN.CENTER)

# Footer credit card
add_rect(s, 3.40, 5.80, 6.50, 1.10, CARD)
add_text(s, "PRESTON CUSICK · ANEESH HARWALKAR",
         3.40, 5.95, 6.50, 0.40, size=11, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, "Applied Human Factors · Spring 2026",
         3.40, 6.30, 6.50, 0.35, size=13, color=TEXT, align=PP_ALIGN.CENTER)
add_text(s, "Prof. John L. Neumann, Ph.D.",
         3.40, 6.55, 6.50, 0.30, size=11, color=DIM, italic=True, align=PP_ALIGN.CENTER)

add_pagenum(s, 11)
add_notes(s, """SPEAKER: Preston · 0:25

"The course taught us one tenet above all: error is a system property, not a human shortcoming.

Duolingo's learners aren't failing. The system is measuring the wrong thing, so it's teaching the wrong thing.

We built a piece that measures what matters. Thank you."

[~25 seconds. End at 8:00.]
""")


# ---- Save ----
out = "/tmp/deck_audit/lightning_talk.pptx"
prs.save(out)
print(f"Saved: {out}  ({len(prs.slides)} slides)")
