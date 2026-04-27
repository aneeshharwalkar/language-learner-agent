#!/usr/bin/env python3
"""
Language Learner Agent — Final Research Report Slide Deck Builder
Applied Human Factors, Spring 2026 · Cusick + Harwalkar

Dark theme (Tract palette), 16:9, 34 slides.
Output: docs/language-learner-agent-research-report.pptx

Language target: stakeholder-friendly plain English.
Technical terms kept only where the professor taught them (see docs/CLAUDE.md).
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ===================== DESIGN TOKENS =====================
BG          = RGBColor(0x0b, 0x0d, 0x10)
SURFACE     = RGBColor(0x13, 0x16, 0x1b)
CARD        = RGBColor(0x1a, 0x1e, 0x26)
BORDER      = RGBColor(0x28, 0x2c, 0x35)

BLUE        = RGBColor(0x7b, 0x9f, 0xe0)
PINK        = RGBColor(0xe0, 0x7b, 0xaa)
GREEN       = RGBColor(0x7b, 0xe0, 0xb4)
AMBER       = RGBColor(0xe0, 0xc2, 0x7b)
RED         = RGBColor(0xe0, 0x7b, 0x7b)

TEXT        = RGBColor(0xf0, 0xf2, 0xf5)
TEXT_SOFT   = RGBColor(0xc8, 0xcd, 0xd8)
MUTED       = RGBColor(0x7a, 0x82, 0x90)

FONT = "Inter"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN_L = Inches(0.6)
CONTENT_W = Inches(13.333 - 1.2)

# ===================== HELPERS =====================

def no_shadow(shape):
    sppr = shape._element.spPr
    for effect in sppr.findall(qn('a:effectLst')):
        sppr.remove(effect)
    etree.SubElement(sppr, qn('a:effectLst'))

def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width_emu=None, rounded=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill_color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        if line_width_emu:
            shp.line.width = line_width_emu
    no_shadow(shp)
    if rounded:
        shp.adjustments[0] = 0.07
    return shp

def add_text(slide, left, top, width, height, text, *,
             size=14, color=TEXT, bold=False, italic=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.25,
             font=FONT, letter_spacing=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    lines = text.split('\n') if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic
        if letter_spacing is not None:
            rPr = run._r.get_or_add_rPr()
            rPr.set('spc', str(int(letter_spacing * 100)))
    return tb

# ===================== PAGE CHROME =====================

def add_section_label(slide, text, color=MUTED):
    add_text(slide, MARGIN_L, Inches(0.45), Inches(7.5), Inches(0.3),
             text.upper(), size=9.5, color=color, bold=True, letter_spacing=1.5)
    add_rect(slide, Inches(5.6), Inches(0.56), Inches(7.2), Emu(6350),
             fill_color=BORDER, line_color=None)

def add_title(slide, text, top=Inches(0.9), size=28, color=TEXT):
    add_text(slide, MARGIN_L, top, CONTENT_W, Inches(1.5),
             text, size=size, color=color, bold=False, spacing=1.18)

def add_subtitle(slide, text, top=Inches(1.85), size=15, color=TEXT_SOFT):
    add_text(slide, MARGIN_L, top, CONTENT_W, Inches(0.5),
             text, size=size, color=color, spacing=1.3)

def add_footer(slide, page_num=None, label=None):
    # Label intentionally omitted per feedback; page number only.
    if page_num is not None:
        add_text(slide, Inches(12.3), Inches(7.12), Inches(0.8), Inches(0.3),
                 f"{page_num}", size=9, color=MUTED, align=PP_ALIGN.RIGHT)

# ===================== CARDS =====================

def add_accent_card(slide, left, top, width, height, title, body, accent,
                    title_size=12, body_size=11, body_color=TEXT_SOFT, title_gap=0.55):
    add_rect(slide, left, top, width, height, fill_color=CARD, line_color=BORDER, line_width_emu=Emu(6350))
    add_rect(slide, left, top, Emu(50800), height, fill_color=accent, line_color=None)
    inner_left = left + Inches(0.22)
    inner_w = width - Inches(0.4)
    add_text(slide, inner_left, top + Inches(0.15), inner_w, Inches(0.45),
             title.upper(), size=title_size, color=accent, bold=True, letter_spacing=1.0, spacing=1.15)
    add_text(slide, inner_left, top + Inches(title_gap), inner_w, height - Inches(title_gap + 0.1),
             body, size=body_size, color=body_color, spacing=1.4)

def add_explainer(slide, left, top, width, height, title, body, accent=BLUE):
    """Plain-English explainer panel for a term that needs a gloss."""
    add_rect(slide, left, top, width, height, fill_color=SURFACE, line_color=BORDER, line_width_emu=Emu(6350))
    add_rect(slide, left, top, Emu(50800), height, fill_color=accent, line_color=None)
    add_text(slide, left + Inches(0.3), top + Inches(0.12), width - Inches(0.45), Inches(0.3),
             title.upper(), size=10, color=accent, bold=True, letter_spacing=1.2)
    add_text(slide, left + Inches(0.3), top + Inches(0.42), width - Inches(0.45), height - Inches(0.52),
             body, size=10.5, color=TEXT_SOFT, spacing=1.4)

def add_pull_quote(slide, left, top, width, height, quote, attribution=None, accent=BLUE):
    add_rect(slide, left, top, width, height, fill_color=SURFACE, line_color=None)
    add_rect(slide, left, top, Emu(50800), height, fill_color=accent, line_color=None)
    add_text(slide, left + Inches(0.3), top + Inches(0.2), width - Inches(0.45), height - Inches(0.5),
             quote, size=13, color=TEXT, italic=True, spacing=1.4)
    if attribution:
        add_text(slide, left + Inches(0.3), top + height - Inches(0.35), width - Inches(0.45), Inches(0.3),
                 attribution, size=10, color=MUTED)

def add_placeholder_image(slide, left, top, width, height, caption_lines, accent=GREEN):
    add_rect(slide, left, top, width, height, fill_color=SURFACE, line_color=accent, line_width_emu=Emu(12700))
    add_text(slide, left + Inches(0.3), top + Inches(0.2), width - Inches(0.6), Inches(0.4),
             "SCREENSHOT PLACEHOLDER", size=10, color=accent, bold=True, letter_spacing=1.5)
    add_text(slide, left + Inches(0.3), top + Inches(0.7), width - Inches(0.6), height - Inches(1.0),
             caption_lines, size=11, color=TEXT_SOFT, spacing=1.4)

def slide_blank(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    return slide

# ===================== OPENING (3 slides) =====================

def slide_01_title(prs):
    slide = slide_blank(prs)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), fill_color=BLUE, line_color=None)
    add_text(slide, MARGIN_L, Inches(1.4), Inches(12), Inches(0.3),
             "APPLIED HUMAN FACTORS · SPRING 2026 · FINAL RESEARCH REPORT",
             size=10, color=MUTED, bold=True, letter_spacing=1.8)
    add_text(slide, MARGIN_L, Inches(2.0), Inches(12), Inches(2.8),
             "Deploying Conversational Agents to Bridge the\nLearning Transfer Gap in Duolingo",
             size=40, color=TEXT, bold=True, spacing=1.1)
    add_text(slide, MARGIN_L, Inches(4.5), Inches(12), Inches(0.5),
             "Evaluate  ·  Analyze  ·  Improve", size=18, color=BLUE, italic=True, letter_spacing=2.0)
    add_text(slide, MARGIN_L, Inches(6.5), Inches(12), Inches(0.4),
             "Preston Cusick and Aneesh Harwalkar", size=13, color=TEXT_SOFT)
    add_text(slide, MARGIN_L, Inches(6.85), Inches(12), Inches(0.3),
             "April 24, 2026", size=10, color=MUTED)
    add_rect(slide, Inches(11.5), Inches(6.7), Inches(1.3), Inches(0.6), fill_color=PINK, line_color=None)

def slide_02_agenda(prs):
    slide = slide_blank(prs)
    add_section_label(slide, "AGENDA")
    add_title(slide, "Five phases. One thesis: system failure, not user failure.")

    phases = [
        ("Phase 1",  "SYSTEM & CONTEXT",       "Duolingo as a human machine system at 500M learner scale.", RED),
        ("Phase 2",  "TASK & USER ANALYSIS",   "Cognitive, perceptual, motor, and workload demands on the learner.", AMBER),
        ("Phase 3",  "DATA & METHODOLOGY",     "A 15 source research corpus and a proposed 24 week experiment.", BLUE),
        ("Phase 4",  "HF ANALYSIS & FINDINGS", "An HFACS walk across all four Swiss Cheese layers.", PINK),
        ("Phase 5",  "REDESIGN",               "Three pillars, a working prototype, and a multi agent future.", GREEN),
    ]
    card_top = Inches(2.5)
    card_h = Inches(3.4)
    total_w = Inches(12.1)
    gap = Inches(0.15)
    card_w = Emu(int(total_w.emu / 5 - gap.emu * 4 / 5))
    x = MARGIN_L
    for num, name, body, color in phases:
        add_rect(slide, x, card_top, card_w, card_h, fill_color=CARD, line_color=BORDER, line_width_emu=Emu(6350))
        add_rect(slide, x, card_top, card_w, Emu(38100), fill_color=color, line_color=None)
        add_text(slide, x + Inches(0.2), card_top + Inches(0.2), card_w - Inches(0.4), Inches(0.4),
                 num.upper(), size=10, color=color, bold=True, letter_spacing=1.5)
        add_text(slide, x + Inches(0.2), card_top + Inches(0.6), card_w - Inches(0.4), Inches(0.8),
                 name, size=13, color=TEXT, bold=True, spacing=1.2)
        add_text(slide, x + Inches(0.2), card_top + Inches(1.6), card_w - Inches(0.4), card_h - Inches(1.8),
                 body, size=11, color=TEXT_SOFT, spacing=1.45)
        x += card_w + gap
    add_text(slide, MARGIN_L, Inches(6.1), CONTENT_W, Inches(0.5),
             "Macro arc: Evaluate, Analyze, Improve.",
             size=12, color=BLUE, italic=True, align=PP_ALIGN.CENTER)
    add_footer(slide, page_num=2)

def slide_03_hook(prs):
    slide = slide_blank(prs)
    add_section_label(slide, "WHY THIS PROJECT")
    add_title(slide, "Months of streaks. And still, no conversation.")
    add_subtitle(slide, "A pattern we have all seen, and the reason this report exists.")

    add_rect(slide, MARGIN_L, Inches(2.9), CONTENT_W, Inches(1.7),
             fill_color=CARD, line_color=BORDER, line_width_emu=Emu(6350))
    add_rect(slide, MARGIN_L, Inches(2.9), Emu(50800), Inches(1.7),
             fill_color=AMBER, line_color=None)
    add_text(slide, Inches(0.95), Inches(3.05), Inches(12), Inches(0.4),
             "THE PATTERN", size=11, color=AMBER, bold=True, letter_spacing=1.2)
    add_text(slide, Inches(0.95), Inches(3.5), Inches(12), Inches(1.0),
             "A learner opens Duolingo every day. A long streak. Thousands of XP. The bird is happy. "
             "Then the moment arrives. A trip abroad, a new coworker, a family visit. "
             "The learner opens their mouth and nothing comes out.\n"
             "This is the transfer gap. The app is the system. The moment on the sidewalk is where transfer either happened or it did not.",
             size=12, color=TEXT_SOFT, spacing=1.45)

    y = Inches(4.85)
    h = Inches(1.9)
    w = Inches(6.05)
    gap = Inches(0.25)
    add_accent_card(slide, MARGIN_L, y, w, h,
                    "WHAT THIS IS NOT",
                    "Not a story about lazy learners. Not a story about bad apps.\n"
                    "The learner did what the system asked them to do. The streak is real. The XP is real.",
                    accent=RED, body_size=11)
    add_accent_card(slide, MARGIN_L + w + gap, y, w, h,
                    "WHAT THIS IS",
                    "A Human Factors investigation of why the system does not produce the outcome its learners came for, "
                    "and a prototype redesign that puts the right signal in the loop.",
                    accent=GREEN, body_size=11)

    add_footer(slide, page_num=3)

# ===================== PHASE 1 · SYSTEM & CONTEXT (6 slides) =====================

def slide_04_system(prs):
    slide = slide_blank(prs)
    add_section_label(slide, "PHASE 1 · SYSTEM & CONTEXT", color=RED)
    add_title(slide, "Duolingo is a 500 million learner human machine system.")
    add_subtitle(slide, "The people who build it, the people who use it, and the environment they use it in.")

    col_top = Inches(3.0)
    col_h = Inches(3.3)
    col_w = Inches(3.95)
    gap = Inches(0.15)
    cols = [
        ("THE PEOPLE WHO BUILD IT", BLUE,
         "• Content and curriculum designers who write lessons.\n"
         "• The engineering team behind Birdbrain, Duolingo's in house algorithm that decides which exercise you see next.\n"
         "• Growth and engagement product teams who manage streaks, notifications, and daily return."),
        ("THE PEOPLE WHO USE IT", AMBER,
         "• 500M+ monthly active learners, mostly adults, mostly self directed.\n"
         "• Mobile first. No classroom and no in person tutor.\n"
         "• Motivated by visible progress: streak count, XP, crown level.\n"
         "• Stated goal: hold a real conversation.\n"
         "• Measured behavior: daily login."),
        ("THE ENVIRONMENT", GREEN,
         "• Short sessions under 15 minutes on a phone.\n"
         "• In between moments: commute, couch, waiting room.\n"
         "• A dense gamified interface.\n"
         "• No other person is in the loop.\n"
         "• The same visual context every session."),
    ]
    x = MARGIN_L
    for title, color, body in cols:
        add_accent_card(slide, x, col_top, col_w, col_h, title=title, body=body, accent=color)
        x += col_w + gap
    add_footer(slide, page_num=4)

def slide_05_stressors(prs):
    slide = slide_blank(prs)
    add_section_label(slide, "PHASE 1 · SYSTEM & CONTEXT", color=RED)
    add_title(slide, "Every human machine system sits under three stressor layers.")
    add_subtitle(slide, "Duolingo's are unusually pronounced, and unusually aligned with how the business makes money.")

    layers = [
        ("INTERNAL STRESSORS", RED,
         "The learner's own mind at work. A feeling of fluency that comes from seeing a word again. "
         "Attention pulled toward XP counters and streak flames instead of the lesson itself."),
        ("EXTERNAL STRESSORS", AMBER,
         "Streaks, XP, leaderboards, crown progression, and push notifications engineered to create daily return pressure. "
         "The app rewards showing up far more than it rewards learning."),
        ("ORGANIZATIONAL STRESSORS", BLUE,
         "Daily Active Users is the primary business metric. Christensen's Innovator's Dilemma at work: a profitable item "
         "based product is hard to leave behind, even when the better version is conversational. Outcome research only "
         "covers reading and listening."),
    ]
    card_top = Inches(2.9)
    card_h = Inches(3.0)
    card_w = Inches(3.95)
    gap = Inches(0.15)
    x = MARGIN_L
    for title, color, body in layers:
        add_accent_card(slide, x, card_top, card_w, card_h, title=title, body=body, accent=color, body_size=11)
        x += card_w + gap
    add_pull_quote(slide, MARGIN_L, Inches(6.15), CONTENT_W, Inches(0.7),
                   "These three layers set the stage. The rest of the report traces how they line up into an incident.",
                   accent=BLUE)
    add_footer(slide, page_num=5)

def slide_06_historical(prs):
    slide = slide_blank(prs)
    add_section_label(slide, "PHASE 1 · SYSTEM & CONTEXT", color=RED)
    add_title(slide, "The constraint used to be technology. It isn't anymore.")
    add_subtitle(slide, "The gap that the 2012 design cannot close, and the 2026 technology that now can.")

    y = Inches(3.3)
    node_w = Inches(3.95)
    node_h = Inches(2.6)
    gap = Inches(0.15)
    timeline = [
        ("2012", "DUOLINGO LAUNCHES",
         "Gamified and item based. Consumer language apps don't have conversational AI. The technology doesn't exist yet.", BLUE),
        ("2023", "GENERATIVE AI ARRIVES",
         "Conversational agents with distinct personas become practical at consumer scale. The technology barrier falls.", AMBER),
        ("2026", "THE GAP PERSISTS",
         "Duolingo's core architecture is unchanged. Jiang et al. (2021) still describes the ceiling. The constraint is now design, not compute.", PINK),
    ]
    x = MARGIN_L
    for year, label, body, color in timeline:
        add_rect(slide, x, y, node_w, node_h, fill_color=CARD, line_color=BORDER, line_width_emu=Emu(6350))
        add_rect(slide, x, y, node_w, Emu(38100), fill_color=color, line_color=None)
        add_text(slide, x + Inches(0.3), y + Inches(0.25), node_w - Inches(0.5), Inches(0.8),
                 year, size=34, color=color, bold=True)
        add_text(slide, x + Inches(0.3), y + Inches(1.0), node_w - Inches(0.5), Inches(0.35),
                 label, size=10, color=MUTED, bold=True, letter_spacing=1.2)
        add_text(slide, x + Inches(0.3), y + Inches(1.35), node_w - Inches(0.5), node_h - Inches(1.5),
                 body, size=11, color=TEXT_SOFT, spacing=1.4)
        x += node_w + gap
    add_footer(slide, page_num=6)

def slide_07_tested_failures(prs):
    slide = slide_blank(prs)
    add_section_label(slide, "PHASE 1 · SYSTEM & CONTEXT", color=RED)
    add_title(slide, "What's actually been tested.", size=26)
    add_subtitle(slide, "The published evidence on Duolingo and mobile language apps, in plain terms.")

    y = Inches(2.7)
    row_h = Inches(0.88)
    gap = Inches(0.1)
    rows = [
        ("JIANG ET AL. (2021)",
         "n = 225 adults, 4 semesters of Duolingo study. Reading reached Intermediate Low. Listening reached Novice High. Speaking and writing were never assessed.",
         RED),
        ("BICKNELL & BRUST (2020) · DUOLINGO'S OWN DESCRIPTION",
         "Birdbrain, Duolingo's in house AI, is optimized to pick the next exercise you are likely to get right. Not calibrated to speaking.",
         AMBER),
        ("MIHAYLOVA ET AL. (2022) · META ANALYSIS OF 23 STUDIES",
         "Mobile apps show large short term gains. But only 3 of 23 studies measured retention after the intervention. Those 3 showed small effects.",
         BLUE),
        ("KRISTENSEN ET AL. (2024) · 8 WEEK APP STUDY",
         "The rate at which learners repeated the same mistakes stayed stable over 8 weeks. More app time did not reduce the pattern. (Caveat: Norwegian second graders.)",
         PINK),
    ]
    for i, (lbl, body, color) in enumerate(rows):
        yy = y + i * (row_h + gap)
        add_rect(slide, MARGIN_L, yy, CONTENT_W, row_h, fill_color=CARD, line_color=BORDER, line_width_emu=Emu(6350))
        add_rect(slide, MARGIN_L, yy, Emu(50800), row_h, fill_color=color, line_color=None)
        add_text(slide, Inches(0.95), yy + Inches(0.15), Inches(5.0), Inches(0.35),
                 lbl, size=10.5, color=color, bold=True, letter_spacing=1.0)
        add_text(slide, Inches(0.95), yy + Inches(0.48), Inches(12), row_h - Inches(0.55),
                 body, size=11, color=TEXT_SOFT, spacing=1.35)

    add_explainer(slide, MARGIN_L, Inches(6.65), CONTENT_W, Inches(0.75),
                  "ACTFL, IN ONE LINE",
                  "The national scale universities use to rate how well someone can use a foreign language. Intermediate Low means simple daily transactions. Real conversation sits two or three levels higher.",
                  accent=BLUE)
    add_footer(slide, page_num=7)

def slide_08_why_matters(prs):
    slide = slide_blank(prs)
    add_section_label(slide, "PHASE 1 · SYSTEM & CONTEXT", color=RED)
    add_title(slide, "A 500 million learner system that doesn't measure the skill its users came for.")

    y = Inches(2.7)
    h = Inches(2.2)
    w = Inches(3.95)
    gap = Inches(0.15)
    stats = [
        ("500M+", "MONTHLY ACTIVE LEARNERS", "One of the largest learning systems on Earth.", BLUE),
        ("0",     "SPEAKING ASSESSMENTS", "In the peer reviewed Duolingo outcome research we surveyed.", RED),
        ("~14%",  "OF US ADULTS", "Report learning a second language. Duolingo dominates the tool choice.", AMBER),
    ]
    x = MARGIN_L
    for val, lbl, desc, color in stats:
        add_rect(slide, x, y, w, h, fill_color=CARD, line_color=BORDER, line_width_emu=Emu(6350))
        add_rect(slide, x, y, w, Emu(38100), fill_color=color, line_color=None)
        add_text(slide, x + Inches(0.25), y + Inches(0.3), w - Inches(0.4), Inches(0.9),
                 val, size=36, color=TEXT, bold=True)
        add_text(slide, x + Inches(0.25), y + Inches(1.15), w - Inches(0.4), Inches(0.3),
                 lbl, size=10, color=MUTED, letter_spacing=1.2, bold=True)
        add_text(slide, x + Inches(0.25), y + Inches(1.45), w - Inches(0.4), Inches(0.6),
                 desc, size=11, color=TEXT_SOFT, spacing=1.35)
        x += w + gap

    add_pull_quote(slide, MARGIN_L, Inches(5.3), CONTENT_W, Inches(1.5),
                   "When the world's largest learning system doesn't measure the skill its users say they want, the result isn't individual under performance. It is a system wide mismatch between what gets measured and what learners actually need, and that is what makes this a Human Factors problem worth solving.",
                   accent=BLUE)
    add_footer(slide, page_num=8)

def slide_09_phase1_close(prs):
    slide = slide_blank(prs)
    add_section_label(slide, "PHASE 1 · SYSTEM & CONTEXT", color=RED)
    add_title(slide, "If this were an aviation incident, we would call it a system failure.")
    add_subtitle(slide, "What we know now, and what the next generation of technology makes possible.")

    y = Inches(3.0)
    h = Inches(3.0)
    w = Inches(6.05)
    gap = Inches(0.25)
    cols = [
        ("WHAT WE KNOW NOW", AMBER,
         "• A measured ceiling at Intermediate Low reading and Novice High listening.\n"
         "• Incentives built around daily return, not real transfer.\n"
         "• Hundreds of millions of engaged learners with rising streaks.\n"
         "• No measured speaking gains in the published research.\n\n"
         "The holes are already open. Phase 4 shows how they line up."),
        ("WHAT COMES NEXT", GREEN,
         "• Generative AI is the first technology positioned to close this gap. Conversational, dialect aware, real time.\n\n"
         "• But only if we name the right enemy first: the business metric that excludes the goal.\n\n"
         "• Phase 5 is our architectural response."),
    ]
    x = MARGIN_L
    for title, color, body in cols:
        add_accent_card(slide, x, y, w, h, title=title, body=body, accent=color, body_size=11)
        x += w + gap

    add_text(slide, MARGIN_L, Inches(6.4), CONTENT_W, Inches(0.5),
             "System failure, not user failure.",
             size=18, color=RED, italic=True, align=PP_ALIGN.CENTER)
    add_footer(slide, page_num=9)

# ===================== PHASE 2 · TASK & USER ANALYSIS (4 slides) =====================

def slide_10_persona(prs):
    slide = slide_blank(prs)
    add_section_label(slide, "PHASE 2 · TASK & USER ANALYSIS", color=AMBER)
    add_title(slide, "Who the learner actually is.")
    add_subtitle(slide, "A portrait of the adult self directed learner the app has to work for.")

    y = Inches(2.8)
    add_rect(slide, MARGIN_L, y, CONTENT_W, Inches(3.6), fill_color=CARD, line_color=BORDER, line_width_emu=Emu(6350))
    add_rect(slide, MARGIN_L, y, Inches(0.06), Inches(3.6), fill_color=AMBER, line_color=None)
    add_text(slide, Inches(0.95), y + Inches(0.2), Inches(11.5), Inches(0.4),
             "PERSONA · THE SELF DIRECTED ADULT LEARNER",
             size=11, color=AMBER, bold=True, letter_spacing=1.5)
    left_body = (
        "• Adult. Self directed. No formal classroom partner.\n"
        "• Mobile first. Daily session under 15 minutes. No one else in the loop.\n"
        "• Motivated by visible progress: streak count, XP, crown level.\n"
        "• Learning for a reason: a trip, a job, a family visit."
    )
    right_body = (
        "• Stated goal: hold a real conversation in the target language.\n"
        "• Actual measured behavior: daily login and XP.\n"
        "• Never speaks inside the app.\n"
        "• Believes the streak is proof of progress."
    )
    add_text(slide, Inches(0.95), y + Inches(0.8), Inches(5.7), Inches(2.6),
             left_body, size=12, color=TEXT_SOFT, spacing=1.55)
    add_text(slide, Inches(6.8), y + Inches(0.8), Inches(5.7), Inches(2.6),
             right_body, size=12, color=TEXT_SOFT, spacing=1.55)

    add_pull_quote(slide, MARGIN_L, Inches(6.6), CONTENT_W, Inches(0.65),
                   "The learner wants to speak. The system rewards showing up.",
                   accent=AMBER)
    add_footer(slide, page_num=10)

def slide_11_recognition_vs_production(prs):
    slide = slide_blank(prs)
    add_section_label(slide, "PHASE 2 · TASK & USER ANALYSIS", color=AMBER)
    add_title(slide, "What the app measures vs. what the learner signed up for.")
    add_subtitle(slide, "Two different skills. The app exercises one. The learner's goal lives in the other.")

    y = Inches(2.9)
    h = Inches(3.3)
    w = Inches(6.05)
    gap = Inches(0.25)

    left_body = (
        "• Multiple choice translation.\n"
        "• Word bank sentence assembly.\n"
        "• Listening practice with clean machine voice clips.\n"
        "• Tap to match vocabulary.\n"
        "• Fill in the blank grammar.\n\n"
        "Demand on the learner:\n"
        "• Perceptual: scan the screen and pick.\n"
        "• Motor: tap, swipe, type.\n"
        "• Cognitive: pick from a small set of options."
    )
    right_body = (
        "• Speaking in real time, without prompts.\n"
        "• Picking the right words for the situation.\n"
        "• Understanding messy, accented speech.\n"
        "• Recovering after a misunderstanding.\n"
        "• Reading the formality of the room.\n\n"
        "Demand on the learner:\n"
        "• Perceptual: hearing messy, accented speech.\n"
        "• Motor: speaking out loud.\n"
        "• Cognitive: build a sentence on the fly, under time pressure."
    )
    add_accent_card(slide, MARGIN_L, y, w, h,
                    "RECOGNITION · WHAT DUOLINGO MEASURES",
                    left_body, accent=BLUE, body_size=11)
    add_accent_card(slide, MARGIN_L + w + gap, y, w, h,
                    "PRODUCTION · WHAT THE LEARNER SIGNED UP FOR",
                    right_body, accent=PINK, body_size=11)

    add_footer(slide, page_num=11)

def slide_12_hf_lenses(prs):
    slide = slide_blank(prs)
    add_section_label(slide, "PHASE 2 · TASK & USER ANALYSIS", color=AMBER)
    add_title(slide, "Cognitive, perceptual, motor, and workload demands.")
    add_subtitle(slide, "Where Duolingo over loads and under loads each kind of demand on the learner.")

    y = Inches(2.8)
    h = Inches(1.9)
    w = Inches(6.05)
    gap_x = Inches(0.25)
    gap_y = Inches(0.2)

    quadrants = [
        ("COGNITIVE · UNDER LOADED", BLUE,
         "No problem solving, no situational reasoning, no sentence building on the fly. The task is pick from a small menu of choices."),
        ("PERCEPTUAL · OVER LOADED (ON THE WRONG CHANNEL)", AMBER,
         "The visual channel is crowded with UI ornaments (streak flame, XP counter, badges). Meanwhile the audio channel is under loaded: clean machine voice only, no accented or messy speech to train on."),
        ("MOTOR · UNDER LOADED", GREEN,
         "Tap and swipe only. The muscles the learner needs for speaking never get exercised inside the app."),
        ("WORKLOAD · UNDER LOADED", PINK,
         "Kept artificially low. Tuned to protect daily retention. Friction is the enemy of showing up, so the system never stretches the learner."),
    ]
    for i, (title, color, body) in enumerate(quadrants):
        row = i // 2
        col = i % 2
        x = MARGIN_L + col * (w + gap_x)
        yy = y + row * (h + gap_y)
        add_accent_card(slide, x, yy, w, h, title, body, accent=color, body_size=11)
    add_footer(slide, page_num=12)

def slide_13_ideal_vs_failure(prs):
    slide = slide_blank(prs)
    add_section_label(slide, "PHASE 2 · TASK & USER ANALYSIS", color=AMBER)
    add_title(slide, "Ideal state versus failure state.")
    add_subtitle(slide, "What the learner signed up for, and where the system actually lands.")

    y = Inches(2.9)
    h = Inches(3.3)
    w = Inches(6.05)
    gap = Inches(0.25)

    ideal_body = (
        "• Hold a five minute conversation in the target language.\n"
        "• With a native speaker, in country, in context.\n"
        "• Recover from a misunderstanding without panic.\n"
        "• Pick the right words for a market, a restaurant, or a workplace.\n"
        "• Understand accented, spontaneous speech.\n\n"
        "This is what the user signed up for."
    )
    now_body = (
        "• Complete daily streak. Earn XP.\n"
        "• Recognize vocabulary in isolation.\n"
        "• Feel confident, without any way to check whether the confidence is earned (Kornell, 2008).\n"
        "• Cannot speak spontaneously under real conditions.\n"
        "• Cannot pick between regional accents of the same language.\n\n"
        "This is where the system actually lands."
    )
    add_accent_card(slide, MARGIN_L, y, w, h,
                    "IDEAL STATE", ideal_body, accent=GREEN, body_size=11)
    add_accent_card(slide, MARGIN_L + w + gap, y, w, h,
                    "FAILURE STATE (NOW)", now_body, accent=RED, body_size=11)

    add_text(slide, MARGIN_L, Inches(6.55), CONTENT_W, Inches(0.4),
             "Phase 3 is built to trace the gap between these two columns.",
             size=12, color=TEXT_SOFT, italic=True, align=PP_ALIGN.CENTER)
    add_footer(slide, page_num=13)

# ===================== PHASE 3 · DATA & METHODOLOGY (5 slides) =====================

def slide_14_approach(prs):
    slide = slide_blank(prs)
    add_section_label(slide, "PHASE 3 · DATA & METHODOLOGY", color=BLUE)
    add_title(slide, "Two weeks, fifteen sources, one prototype, one proposed experiment.")
    add_subtitle(slide, "Research driven and build driven. Transparent about what it can and cannot prove.")

    y = Inches(2.9)
    h = Inches(3.3)
    w = Inches(6.05)
    gap = Inches(0.25)

    approach_body = (
        "• A 15 source peer reviewed research corpus.\n"
        "• Two meta analyses (Bibauw 2022, Mayer 2021) carry the intervention claims.\n"
        "• A working conversational agent, built in the repo.\n"
        "• A proposed 24 week randomized controlled trial for a follow on study.\n"
        "• No primary data was collected in this window."
    )
    benefits_body = (
        "• Synthesis leverages a far larger evidence base than a two week primary study could ever produce.\n"
        "• The prototype constrains the redesign space to what is actually buildable.\n"
        "• Every claim in this deck ties back to a source file in the repo.\n"
        "• The proposed experiment closes what this window could not."
    )
    add_accent_card(slide, MARGIN_L, y, w, h,
                    "APPROACH", approach_body, accent=BLUE, body_size=13)
    add_accent_card(slide, MARGIN_L + w + gap, y, w, h,
                    "BENEFITS", benefits_body, accent=GREEN, body_size=13)

    add_footer(slide, page_num=14)

def slide_15_build_process(prs):
    slide = slide_blank(prs)
    add_section_label(slide, "PHASE 3 · DATA & METHODOLOGY", color=BLUE)
    add_title(slide, "How we built it: five scoped agents, one coordinated project.")
    add_subtitle(slide, "We used Claude Code to run five sub agents, each with a single role and a single directory.")

    y = Inches(2.85)
    h = Inches(3.3)
    total_w = Inches(12.1)
    gap = Inches(0.1)
    card_w = Emu(int(total_w.emu / 5 - gap.emu * 4 / 5))
    x = MARGIN_L

    agents = [
        ("RESEARCH", RED,
         "Assembled the 15 source corpus. One note per paper. Direct quotes, sample sizes, effect sizes. Flagged paraphrases."),
        ("HF ANALYSIS", AMBER,
         "Produced the HFACS walk across all four Swiss Cheese levels. Mapped every claim to a source and a course concept."),
        ("FRONTEND", BLUE,
         "Built the chat UI in plain HTML, CSS, and JavaScript. Dialect pill, scenario pill, severity badges, translation toggle."),
        ("BACKEND", PINK,
         "Built the Flask server and the agent. Scenario openers, dialect aware prompts, severity graded corrections."),
        ("QA", GREEN,
         "Reviewed every deliverable against the research files and the API contract. Produced the blocker, warning, and note report that drove the final polish."),
    ]
    for name, color, body in agents:
        add_rect(slide, x, y, card_w, h, fill_color=CARD, line_color=BORDER, line_width_emu=Emu(6350))
        add_rect(slide, x, y, card_w, Emu(38100), fill_color=color, line_color=None)
        add_text(slide, x + Inches(0.2), y + Inches(0.2), card_w - Inches(0.4), Inches(0.4),
                 name, size=11, color=color, bold=True, letter_spacing=1.2)
        add_text(slide, x + Inches(0.2), y + Inches(0.7), card_w - Inches(0.4), h - Inches(0.8),
                 body, size=10.5, color=TEXT_SOFT, spacing=1.45)
        x += card_w + gap

    add_pull_quote(slide, MARGIN_L, Inches(6.35), CONTENT_W, Inches(0.75),
                   "One directory per agent. One instruction file per agent. That is how we kept the work coordinated and every deliverable traceable to an owner.",
                   accent=BLUE)
    add_footer(slide, page_num=15)

def slide_16_corpus(prs):
    slide = slide_blank(prs)
    add_section_label(slide, "PHASE 3 · DATA & METHODOLOGY", color=BLUE)
    add_title(slide, "Corpus: 15 peer reviewed sources, grouped by function.")
    add_subtitle(slide, "Two meta analyses carry the intervention claims. Full source files live in the repo.")

    y = Inches(2.7)
    h = Inches(3.0)
    w = Inches(2.95)
    gap = Inches(0.15)
    groups = [
        ("DUOLINGO & MOBILE APPS", RED,
         "• Jiang et al. (2021). ACTFL ceiling study.\n"
         "• Mihaylova et al. (2022). Mobile app meta analysis.\n"
         "• Kristensen et al. (2024). Error repetition.\n"
         "• Bicknell (2020). Birdbrain (Duolingo's own description)."),
        ("MEMORY & CONTEXT", AMBER,
         "• Baddeley (2000). Working Memory and the Episodic Buffer.\n"
         "• Essoe et al. (2022). Context reinstatement.\n"
         "• Godwin-Jones (2018). Isolated drill problem.\n"
         "• Kornell (2008). Why learners can't tell what's working."),
        ("INTERVENTION EVIDENCE", GREEN,
         "• Bibauw et al. (2022). Dialogue practice (d = 0.58).\n"
         "• Mayer (2021). Conversational style (d = 1.3).\n"
         "• Geeslin and Long (2014). Regional language variation.\n"
         "• Wang (2024). GenAI chatbot (affect only)."),
        ("SUPERVISION & TRANSPARENCY", BLUE,
         "• Cannon-Bowers et al. (1998). TADMUS framework.\n"
         "• Loft et al. (2021). Transparency as a system variable (course source)."),
    ]
    x = MARGIN_L
    for title, color, body in groups:
        add_accent_card(slide, x, y, w, h, title, body, accent=color, body_size=10.5, title_gap=0.6)
        x += w + gap

    add_explainer(slide, MARGIN_L, Inches(5.95), CONTENT_W, Inches(1.05),
                  "HOW TO READ EFFECT SIZES (d)",
                  "A standard way to measure how big a gap is between two groups. d = 0.2 is small, d = 0.5 is medium, d = 0.8 is large. Anything above 1.0 is very large for education research. You will see d = 0.58 and d = 1.3 later in this report.",
                  accent=GREEN)
    add_footer(slide, page_num=16)

def slide_17_rct(prs):
    slide = slide_blank(prs)
    add_section_label(slide, "PHASE 3 · DATA & METHODOLOGY", color=BLUE)
    add_title(slide, "Proposed experiment: what a real test would manipulate and measure.")
    add_subtitle(slide, "A 24 week randomized controlled trial. Retention is a co-primary outcome, not a follow up.")

    y = Inches(2.5)
    row_h = Inches(0.62)
    gap = Inches(0.08)
    rows = [
        ("MAIN EFFECT 1 · AGENT TYPE",
         "Four levels: dialect specific (Carlos or Elena), generic chatbot, embodied avatar (per Wang 2024), Duolingo only control.", PINK),
        ("MAIN EFFECT 2 · TIME",
         "Four measurement points: week 0 (baseline), week 12 (post), week 24 (durable), and a 6 month follow up.", GREEN),
        ("INTERACTION · AGENT TYPE × DIALECT MATCH",
         "Does a Mexican Spanish learner gain more from Carlos than from Elena? Does any dialect agent beat a generic one?", AMBER),
        ("INTERACTION · AGENT TYPE × TIME (RETENTION)",
         "Do gains from the dialect agent decay more slowly than gains from the generic agent or the control? This is the transfer test.", BLUE),
        ("WHAT WE WOULD MEASURE",
         "ACTFL Oral Proficiency Interview, scored by a blinded rater. The standard live speaking test universities use.", TEXT_SOFT),
        ("WHAT WE'D CONTROL FOR",
         "Prior proficiency. Baseline streak length. Daily engagement time. Age. First language. Prior motivation. Intended use context.", RED),
        ("POST HOC ANALYSES",
         "Sub scales of speaking (pragmatics, register, repair). Self reported confidence. App engagement during the trial.", MUTED),
    ]
    for i, (lbl, body, color) in enumerate(rows):
        yy = y + i * (row_h + gap)
        add_rect(slide, MARGIN_L, yy, CONTENT_W, row_h, fill_color=CARD, line_color=BORDER, line_width_emu=Emu(6350))
        add_rect(slide, MARGIN_L, yy, Emu(50800), row_h, fill_color=color, line_color=None)
        add_text(slide, Inches(0.9), yy + Inches(0.1), Inches(4.3), row_h - Inches(0.2),
                 lbl, size=10, color=color, bold=True, letter_spacing=1.0)
        add_text(slide, Inches(5.4), yy + Inches(0.1), Inches(7.4), row_h - Inches(0.2),
                 body, size=10.5, color=TEXT_SOFT, spacing=1.3)
    add_footer(slide, page_num=17)

def slide_18_threats(prs):
    slide = slide_blank(prs)
    add_section_label(slide, "PHASE 3 · DATA & METHODOLOGY", color=BLUE)
    add_title(slide, "Threats to validity, and how we would address each one.")

    y = Inches(2.4)
    row_h = Inches(0.7)
    gap = Inches(0.1)
    threats = [
        ("NO PRIMARY DATA", RED,
         "Two week build window. Addressed by relying on two meta analyses (Bibauw, Mayer) and a proposed future experiment."),
        ("NO A/B OF THE PROTOTYPE", RED,
         "A proof of concept only. Addressed by the proposed 24 week study with random assignment."),
        ("SELF SELECTION ON DIALECT", AMBER,
         "Dialect is a product input: learners pick based on where they're going or who they talk to. To isolate the dialect effect, the future RCT randomly assigns dialect and stratifies by self reported intended use context."),
        ("WHY 24 WEEKS, NOT WANG'S 12", AMBER,
         "Wang (2024) ran 12 weeks and found no speaking gain. We extend the window to 24 weeks plus a 6 month follow up so retention has room to show up in the data."),
        ("KRISTENSEN (2024) GENERALIZABILITY", AMBER,
         "That study sampled Norwegian second graders. We use it only for the mechanism it shows, with the caveat stated."),
        ("SHORT TERM BIAS IN THE APP LITERATURE", BLUE,
         "Mihaylova (2022) reports 78% of mobile app studies had high risk of bias and only 22% used randomization. We lean on meta analytic effects, not single study claims."),
    ]
    for i, (lbl, color, body) in enumerate(threats):
        yy = y + i * (row_h + gap)
        add_rect(slide, MARGIN_L, yy, CONTENT_W, row_h, fill_color=CARD, line_color=BORDER, line_width_emu=Emu(6350))
        add_rect(slide, MARGIN_L, yy, Emu(50800), row_h, fill_color=color, line_color=None)
        add_text(slide, Inches(0.9), yy + Inches(0.1), Inches(4.0), row_h - Inches(0.2),
                 lbl, size=10.5, color=color, bold=True, letter_spacing=1.0)
        add_text(slide, Inches(5.1), yy + Inches(0.1), Inches(7.7), row_h - Inches(0.2),
                 body, size=10.5, color=TEXT_SOFT, spacing=1.3)
    add_footer(slide, page_num=18)

# ===================== PHASE 4 · HF ANALYSIS & FINDINGS (8 slides) =====================

def slide_19_hfacs_method(prs):
    slide = slide_blank(prs)
    add_section_label(slide, "PHASE 4 · HF ANALYSIS & FINDINGS", color=PINK)
    add_title(slide, "Reason's Swiss Cheese Model. We walk it backward, not downward.")
    add_subtitle(slide, "HFACS is the method used to investigate aviation and industrial incidents. Same method, new domain.")

    x_slice = MARGIN_L
    w_slice = Inches(6.0)
    slice_h = Inches(0.7)
    slice_gap = Inches(0.15)
    slice_top = Inches(2.9)
    layers = [
        ("L4 · ORGANIZATIONAL INFLUENCES", BLUE,  "Culture, policy, business metrics, resources"),
        ("L3 · UNSAFE SUPERVISION",        PINK,  "Oversight, planning, algorithmic control"),
        ("L2 · PRECONDITIONS",             AMBER, "Environment, mental state, context"),
        ("L1 · UNSAFE ACTS",               RED,   "Errors and violations at the operator level"),
    ]
    for i, (title, color, subtitle) in enumerate(layers):
        yy = slice_top + i * (slice_h + slice_gap)
        add_rect(slide, x_slice, yy, w_slice, slice_h, fill_color=CARD, line_color=BORDER, line_width_emu=Emu(6350), rounded=True)
        add_rect(slide, x_slice, yy, Emu(50800), slice_h, fill_color=color, line_color=None)
        add_text(slide, x_slice + Inches(0.3), yy + Inches(0.12), Inches(5.5), Inches(0.3),
                 title, size=10.5, color=color, bold=True, letter_spacing=1.0)
        add_text(slide, x_slice + Inches(0.3), yy + Inches(0.37), Inches(5.5), Inches(0.3),
                 subtitle, size=9.5, color=MUTED)

    x_text = Inches(6.9)
    w_text = Inches(6.0)
    add_explainer(slide, x_text, Inches(2.9), w_text, Inches(1.3),
                  "THE SWISS CHEESE MODEL, IN PLAIN TERMS",
                  "Every layer of a system has small weaknesses (holes). Most of the time the holes don't line up. When they do, the failure slips all the way through. Reason's model came out of aviation accident investigation.",
                  accent=PINK)

    add_text(slide, x_text, Inches(4.35), w_text, Inches(0.3),
             "THE METHOD", size=10, color=MUTED, bold=True, letter_spacing=1.5)
    add_text(slide, x_text, Inches(4.7), w_text, Inches(2.2),
             "• Identify the observable failure. Learners can't speak after months of app use.\n"
             "• Work backward through all four HFACS layers.\n"
             "• Use evidence to support each layer.\n"
             "• Look for the system breakdown. Avoid blaming the individual.\n"
             "• Propose system level interventions in Phase 5.",
             size=11, color=TEXT_SOFT, spacing=1.5)
    add_text(slide, x_text, Inches(6.85), w_text, Inches(0.3),
             "Course tenet: error is a system property, not a human shortcoming. Exxon Valdez. Sully. Now Duolingo.",
             size=11, color=PINK, italic=True)
    add_footer(slide, page_num=19)

def slide_20_l4_dau(prs):
    slide = slide_blank(prs)
    add_section_label(slide, "PHASE 4 · LEVEL 4 · ORGANIZATIONAL INFLUENCES", color=BLUE)
    add_title(slide, "Daily Active Users is the wrong metric to build around.")
    add_subtitle(slide, "And the business is locked into keeping it. This is the Innovator's Dilemma at work.")

    add_explainer(slide, MARGIN_L, Inches(2.75), CONTENT_W, Inches(1.0),
                  "THE INNOVATOR'S DILEMMA (CHRISTENSEN)",
                  "Established companies are so good at their existing, profitable product that they get blindsided by disruptive alternatives. The best thing for the business in the short term is also what blocks the bigger change.",
                  accent=AMBER)

    y = Inches(4.0)
    h = Inches(1.0)
    gap_y = Inches(0.15)
    rows = [
        ("THE METRIC",
         "Daily Active Users is Duolingo's primary business metric. It shapes streak design, notification timing, exercise sequencing, and what gets measured in outcome research.", BLUE),
        ("THE DILEMMA",
         "A disruptive alternative (dialect specific conversation, adaptive speaking practice) costs more to build and does not move the daily return metric in the short term.", AMBER),
        ("THE CONSEQUENCE",
         "Every layer downstream inherits this optimization. The algorithm, the interface, and the learner all act on what drives daily return. The cheese is pre holed at the top.", RED),
    ]
    for i, (lbl, body, color) in enumerate(rows):
        yy = y + i * (h + gap_y)
        add_rect(slide, MARGIN_L, yy, CONTENT_W, h, fill_color=CARD, line_color=BORDER, line_width_emu=Emu(6350))
        add_rect(slide, MARGIN_L, yy, Emu(50800), h, fill_color=color, line_color=None)
        add_text(slide, Inches(0.95), yy + Inches(0.15), Inches(3.0), Inches(0.3),
                 lbl, size=11, color=color, bold=True, letter_spacing=1.2)
        add_text(slide, Inches(0.95), yy + Inches(0.48), Inches(12), h - Inches(0.55),
                 body, size=11, color=TEXT_SOFT, spacing=1.4)
    add_footer(slide, page_num=20)

def slide_21_l4_personalization(prs):
    slide = slide_blank(prs)
    add_section_label(slide, "PHASE 4 · LEVEL 4 · ORGANIZATIONAL INFLUENCES", color=BLUE)
    add_title(slide, "Conversational style moves learning the most. The business underweights it.")
    add_subtitle(slide, "Mayer's research measures learning outcomes with Cohen's d (effect size). d = 0.8 is already large for education research.")

    y = Inches(2.8)
    h = Inches(1.95)
    w = Inches(6.05)
    gap = Inches(0.25)

    add_rect(slide, MARGIN_L, y, w, h, fill_color=CARD, line_color=BORDER, line_width_emu=Emu(6350))
    add_rect(slide, MARGIN_L, y, w, Emu(38100), fill_color=GREEN, line_color=None)
    add_text(slide, MARGIN_L + Inches(0.25), y + Inches(0.2), w - Inches(0.4), Inches(0.75),
             "median d = 1.3", size=28, color=TEXT, bold=True)
    add_text(slide, MARGIN_L + Inches(0.25), y + Inches(0.95), w - Inches(0.4), Inches(0.35),
             "MAYER (2021) · CONVERSATIONAL STYLE · HOW MUCH LEARNING IMPROVED",
             size=10, color=MUTED, letter_spacing=1.0, bold=True)
    add_text(slide, MARGIN_L + Inches(0.25), y + Inches(1.3), w - Inches(0.4), Inches(0.55),
             "Supported in 10 of 10 experiments. Learners who got a conversational tutor learned much more than learners who got the same content delivered formally.",
             size=10.5, color=TEXT_SOFT, spacing=1.35)

    x2 = MARGIN_L + w + gap
    add_rect(slide, x2, y, w, h, fill_color=CARD, line_color=BORDER, line_width_emu=Emu(6350))
    add_rect(slide, x2, y, w, Emu(38100), fill_color=AMBER, line_color=None)
    add_text(slide, x2 + Inches(0.25), y + Inches(0.2), w - Inches(0.4), Inches(0.75),
             "median d = 0.8", size=28, color=TEXT, bold=True)
    add_text(slide, x2 + Inches(0.25), y + Inches(0.95), w - Inches(0.4), Inches(0.35),
             "MAYER (2021) · NATIVE ACCENTED VOICE · LEARNING IMPROVEMENT",
             size=10, color=MUTED, letter_spacing=1.0, bold=True)
    add_text(slide, x2 + Inches(0.25), y + Inches(1.3), w - Inches(0.4), Inches(0.55),
             "Motivates the proposed voice agent in Phase 5. The current prototype ships conversational style. Voice is future work.",
             size=10.5, color=TEXT_SOFT, spacing=1.35)

    add_rect(slide, MARGIN_L, Inches(5.0), CONTENT_W, Inches(2.0), fill_color=SURFACE, line_color=BORDER, line_width_emu=Emu(6350))
    add_rect(slide, MARGIN_L, Inches(5.0), Emu(50800), Inches(2.0), fill_color=PINK, line_color=None)
    add_text(slide, Inches(0.95), Inches(5.15), Inches(11.5), Inches(0.35),
             "TRANSPARENCY IS A VARIABLE, NOT A CONSTANT (LOFT ET AL., 2021)",
             size=11, color=PINK, bold=True, letter_spacing=1.2)
    add_text(slide, Inches(0.95), Inches(5.55), Inches(11.5), Inches(1.4),
             "Older Human Factors research treated transparency as a fixed feature. The Loft study assigned in this course treats transparency as a design variable that can be tuned. The catch: showing confident projected outcomes can make a user feel justified moving faster, even when the system is wrong.\n"
             "Duolingo's progress UI is exactly this kind of projected outcome. It is highly transparent about XP and streaks, which do not include speaking.",
             size=10.5, color=TEXT_SOFT, spacing=1.4)

    add_footer(slide, page_num=21)

def slide_22_l3(prs):
    slide = slide_blank(prs)
    add_section_label(slide, "PHASE 4 · LEVEL 3 · UNSAFE SUPERVISION", color=PINK)
    add_title(slide, "The supervisor inside the app is optimizing for the wrong signal.")
    add_subtitle(slide, "Birdbrain is the in app algorithm. Signal Detection Theory tells us what it can and cannot catch.")

    add_pull_quote(slide, MARGIN_L, Inches(2.65), Inches(6.0), Inches(1.5),
                   "\"Birdbrain makes an educated guess about whether a learner will get a given exercise right.\"",
                   attribution="Duolingo engineering blog (Bicknell, 2020). Duolingo's own description of what its algorithm optimizes for.",
                   accent=PINK)

    add_explainer(slide, MARGIN_L, Inches(4.3), Inches(6.0), Inches(1.25),
                  "WHAT BIRDBRAIN IS",
                  "Duolingo's in house AI. Every time you open the app, Birdbrain picks the next exercise for you. It is the layer between the content library and the learner.",
                  accent=PINK)

    add_rect(slide, MARGIN_L, Inches(5.7), Inches(6.0), Inches(1.35), fill_color=CARD, line_color=BORDER, line_width_emu=Emu(6350))
    add_rect(slide, MARGIN_L, Inches(5.7), Emu(50800), Inches(1.35), fill_color=RED, line_color=None)
    add_text(slide, Inches(0.95), Inches(5.8), Inches(5.5), Inches(0.3),
             "SIGNAL DETECTION THEORY · WHAT BIRDBRAIN CAN SEE",
             size=10, color=RED, bold=True, letter_spacing=1.2)
    add_text(slide, Inches(0.95), Inches(6.1), Inches(5.5), Inches(0.95),
             "In SDT terms, a detector's sensitivity (d') is how well it separates signal from noise. "
             "Birdbrain's d' for item correctness is high. Its d' for conversational failure is effectively zero. "
             "It cannot detect what it never measures.",
             size=10.5, color=TEXT_SOFT, spacing=1.4)

    add_rect(slide, Inches(6.9), Inches(2.65), Inches(6.0), Inches(4.4), fill_color=CARD, line_color=BORDER, line_width_emu=Emu(6350))
    add_rect(slide, Inches(6.9), Inches(2.65), Emu(50800), Inches(4.4), fill_color=BLUE, line_color=None)
    add_text(slide, Inches(7.25), Inches(2.8), Inches(5.5), Inches(0.35),
             "TADMUS MISMATCH (CANNON-BOWERS, 1998)",
             size=10.5, color=BLUE, bold=True, letter_spacing=1.2)
    add_text(slide, Inches(7.25), Inches(3.2), Inches(5.5), Inches(3.8),
             "TADMUS (Tactical Decision Making Under Stress) is the Navy funded research program we covered in Week 3. "
             "One core principle: operators under pressure need feedback that diagnoses the actual task, not feedback calibrated to a proxy the system happens to be tracking.\n\n"
             "Duolingo's progress display is calibrated to Birdbrain's memorization state. The learner's mental model is "
             "of speaking competence. These are two different things. That is the textbook TADMUS mismatch.\n\n"
             "Godwin-Jones (2018) makes the same point in language teaching terms. Duolingo's spaced drill scheduler optimizes "
             "for memorizing words, not for using them correctly in real conversation.",
             size=10.5, color=TEXT_SOFT, spacing=1.45)

    add_footer(slide, page_num=22)

def slide_23_l2(prs):
    slide = slide_blank(prs)
    add_section_label(slide, "PHASE 4 · LEVEL 2 · PRECONDITIONS", color=AMBER)
    add_title(slide, "The memory system Duolingo never engages.")
    add_subtitle(slide, "Baddeley's Working Memory model has four parts. Duolingo exercises one of them.")

    add_explainer(slide, MARGIN_L, Inches(2.5), CONTENT_W, Inches(0.95),
                  "BADDELEY'S WORKING MEMORY, IN PLAIN TERMS",
                  "The Phonological Loop is the part of working memory that rehearses words in your head (repeating a phone number silently). The Episodic Buffer binds what you are hearing right now to everything you already know: the setting, the situation, the meaning. Vocabulary without the Episodic Buffer does not transfer.",
                  accent=AMBER)

    y = Inches(3.6)
    w = Inches(2.95)
    h = Inches(1.3)
    gap = Inches(0.15)
    components = [
        ("PHONOLOGICAL LOOP", GREEN, "Verbal rehearsal",       "Heavily engaged"),
        ("VISUOSPATIAL SKETCHPAD", MUTED, "Visual imagery",    "Partially engaged"),
        ("CENTRAL EXECUTIVE", MUTED, "Attention and control",  "Partially engaged"),
        ("EPISODIC BUFFER",  RED,   "Binds context and memory","Never engaged"),
    ]
    x = MARGIN_L
    for title, color, func, status in components:
        add_rect(slide, x, y, w, h, fill_color=CARD, line_color=BORDER, line_width_emu=Emu(6350))
        add_rect(slide, x, y, w, Emu(38100), fill_color=color, line_color=None)
        add_text(slide, x + Inches(0.2), y + Inches(0.15), w - Inches(0.4), Inches(0.35),
                 title, size=10, color=color, bold=True, letter_spacing=0.8)
        add_text(slide, x + Inches(0.2), y + Inches(0.55), w - Inches(0.4), Inches(0.3),
                 func, size=10, color=TEXT_SOFT)
        add_text(slide, x + Inches(0.2), y + Inches(0.9), w - Inches(0.4), Inches(0.3),
                 status, size=11, color=color, italic=True, bold=True)
        x += w + gap

    y2 = Inches(5.1)
    add_rect(slide, MARGIN_L, y2, Inches(6.05), Inches(1.9), fill_color=CARD, line_color=BORDER, line_width_emu=Emu(6350))
    add_rect(slide, MARGIN_L, y2, Inches(6.05), Emu(38100), fill_color=GREEN, line_color=None)
    add_text(slide, Inches(0.95), y2 + Inches(0.2), Inches(5.5), Inches(0.35),
             "ESSOE ET AL. (2022) · CONTEXT MATTERS",
             size=10, color=GREEN, bold=True, letter_spacing=1.2)
    add_text(slide, Inches(0.95), y2 + Inches(0.55), Inches(5.5), Inches(0.7),
             "92% versus 76%", size=24, color=TEXT, bold=True)
    add_text(slide, Inches(0.95), y2 + Inches(1.2), Inches(5.5), Inches(0.65),
             "One week retention when vocabulary was studied in distinct contexts versus the same context. Also: 38% fewer cross language mix ups.",
             size=10.5, color=TEXT_SOFT, spacing=1.4)

    x3 = MARGIN_L + Inches(6.05) + Inches(0.25)
    add_pull_quote(slide, x3, y2, Inches(6.05), Inches(1.9),
                   "Duolingo's gamified interface looks the same every session. Context binding never happens. This is why short term memorization gains don't turn into long term conversation skill. The Episodic Buffer simply never gets engaged.",
                   accent=RED)

    add_footer(slide, page_num=23)

def slide_24_l1(prs):
    slide = slide_blank(prs)
    add_section_label(slide, "PHASE 4 · LEVEL 1 · UNSAFE ACTS", color=RED)
    add_title(slide, "The learner acts on the only signal the system displays.")
    add_subtitle(slide, "That signal is the wrong one. The patterns that follow are well studied.")

    y = Inches(2.7)
    h = Inches(1.85)
    w = Inches(6.05)
    gap_x = Inches(0.25)
    gap_y = Inches(0.15)
    panels = [
        ("FALSE ALARM (SDT)", RED,
         "In Signal Detection Theory, a false alarm is when the detector says \"present\" for a signal that isn't actually there. The learner feels confident because XP and streaks are visible. Those signals don't include real speaking ability. The confidence is a false alarm on a signal that was never there."),
        ("KORNELL (2008) · THE CONFIDENCE PROBLEM", PINK,
         "\"Even after testing, and contrary to their own test performance, participants rated massing as more effective than spacing.\" Learners cannot tell what is working. They rely on what the system shows them, and the system is showing them the wrong thing."),
        ("KRISTENSEN (2024) · ERROR REPETITION", AMBER,
         "\"The propensity to repeat mistakes remained stable over time.\" More app time did not reduce the pattern. Caveat: this study sampled Norwegian second graders. We use it for the mechanism, not as evidence about adults specifically."),
        ("PLAN CONTINUATION BIAS", BLUE,
         "A well documented pattern we covered in Week 4 (Air India Express 812). People keep executing a plan even after the evidence says to stop. Learners keep extending streaks past mounting evidence of no real transfer, because the system suppresses the signals that would argue for a change."),
    ]
    for i, (title, color, body) in enumerate(panels):
        row = i // 2
        col = i % 2
        x = MARGIN_L + col * (w + gap_x)
        yy = y + row * (h + gap_y)
        add_accent_card(slide, x, yy, w, h, title, body, accent=color, body_size=10)

    add_explainer(slide, MARGIN_L, Inches(6.55), CONTENT_W, Inches(0.65),
                  "FUNDAMENTAL ATTRIBUTION ERROR (COURSE CONCEPT, WEEK 9)",
                  "Blaming a person's character for behavior that was actually caused by the situation around them. Calling the learner \"lazy\" for missing a streak is textbook FAE when the system never gave them a real signal to act on.",
                  accent=RED)
    add_footer(slide, page_num=24)

def slide_25_capstone(prs):
    slide = slide_blank(prs)
    add_section_label(slide, "PHASE 4 · FAIL POINT MAP", color=PINK)
    add_title(slide, "The holes line up. Name the enemy.")

    x_left = MARGIN_L
    w_left = Inches(6.2)
    top = Inches(2.4)
    slice_h = Inches(0.85)
    slice_gap = Inches(0.1)
    layers = [
        ("L4 · DAILY ACTIVE USER METRIC EXCLUDES SPEAKING", BLUE),
        ("L3 · BIRDBRAIN OPTIMIZES ITEM CORRECTNESS",       PINK),
        ("L2 · EPISODIC BUFFER NEVER ENGAGED",              AMBER),
        ("L1 · THE LEARNER HAS NO REAL SIGNAL",             RED),
    ]
    for i, (label, color) in enumerate(layers):
        yy = top + i * (slice_h + slice_gap)
        add_rect(slide, x_left, yy, w_left, slice_h, fill_color=CARD, line_color=BORDER, line_width_emu=Emu(6350), rounded=True)
        add_rect(slide, x_left, yy, Emu(50800), slice_h, fill_color=color, line_color=None)
        add_text(slide, x_left + Inches(0.3), yy + Inches(0.25), Inches(5.5), Inches(0.4),
                 label, size=11, color=color, bold=True, letter_spacing=1.0, spacing=1.1)

    x_right = Inches(7.1)
    w_right = Inches(5.8)
    enemies = [
        ("01", "THE ISOLATED DRILL",              "Vocabulary taught as isolated items, never embedded in a real situation.", RED),
        ("02", "THE WRONG SIGNAL",                "Confident feedback about XP and streaks. The wrong variable, amplified.", AMBER),
        ("03", "THE RECEPTIVE ONLY EVIDENCE BASE","Published outcome research covers reading and listening. Speaking is unsupervised by design.", PINK),
        ("04", "THE ALIGNED HOLES",               "Fix one layer alone and the cheese still lines up. The redesign has to span layers.", BLUE),
    ]
    card_h = Inches(1.0)
    gap_y = Inches(0.1)
    for i, (num, name, desc, color) in enumerate(enemies):
        yy = top + i * (card_h + gap_y)
        add_rect(slide, x_right, yy, w_right, card_h, fill_color=CARD, line_color=BORDER, line_width_emu=Emu(6350))
        add_rect(slide, x_right, yy, Emu(50800), card_h, fill_color=color, line_color=None)
        add_text(slide, x_right + Inches(0.2), yy + Inches(0.1), Inches(0.5), Inches(0.4),
                 num, size=16, color=color, bold=True)
        add_text(slide, x_right + Inches(0.7), yy + Inches(0.12), w_right - Inches(0.8), Inches(0.4),
                 name, size=11, color=color, bold=True, letter_spacing=1.0)
        add_text(slide, x_right + Inches(0.7), yy + Inches(0.45), w_right - Inches(0.8), Inches(0.5),
                 desc, size=10.5, color=TEXT_SOFT, spacing=1.35)

    add_text(slide, MARGIN_L, Inches(6.75), CONTENT_W, Inches(0.4),
             "Rosling. Tufte. Musk. Name the enemy.",
             size=11, color=MUTED, italic=True, align=PP_ALIGN.CENTER)
    add_footer(slide, page_num=25)

def slide_26_against_grain(prs):
    slide = slide_blank(prs)
    add_section_label(slide, "PHASE 4 · DISCONFIRMING EVIDENCE", color=AMBER)
    add_title(slide, "The finding that complicates our own story.")
    add_subtitle(slide, "Not every AI chatbot has moved the needle on speaking. We have to say so.")

    y = Inches(2.8)
    add_rect(slide, MARGIN_L, y, CONTENT_W, Inches(3.1), fill_color=CARD, line_color=BORDER, line_width_emu=Emu(6350))
    add_rect(slide, MARGIN_L, y, Emu(50800), Inches(3.1), fill_color=AMBER, line_color=None)
    add_text(slide, Inches(0.95), y + Inches(0.2), Inches(12), Inches(0.4),
             "WANG (2024) · AI CHATBOT · 12 WEEK STUDY · n = 99 · INTERMEDIATE ADULT LEARNERS",
             size=10.5, color=AMBER, bold=True, letter_spacing=1.2)

    add_text(slide, Inches(0.95), y + Inches(0.7), Inches(5.8), Inches(0.3),
             "WHAT IT DID MOVE", size=10.5, color=GREEN, bold=True, letter_spacing=1.2)
    add_text(slide, Inches(0.95), y + Inches(1.05), Inches(5.8), Inches(1.8),
             "• Willingness to communicate (less anxiety about speaking)\n"
             "• Self rated ability to hold a conversation\n"
             "• Reported enjoyment of practice",
             size=11.5, color=TEXT_SOFT, spacing=1.5)
    add_text(slide, Inches(7.0), y + Inches(0.7), Inches(5.8), Inches(0.3),
             "WHAT IT DID NOT MOVE", size=10.5, color=RED, bold=True, letter_spacing=1.2)
    add_text(slide, Inches(7.0), y + Inches(1.05), Inches(5.8), Inches(1.8),
             "• Measured speaking performance\n"
             "• Between group differences on objective speaking tests\n"
             "• Transfer to real conversation outside the app",
             size=11.5, color=TEXT_SOFT, spacing=1.5)

    add_pull_quote(slide, MARGIN_L, Inches(6.1), CONTENT_W, Inches(1.0),
                   "Engagement is not transfer. A proof of concept is proof of plausibility, not proof of effect. This is why our proposed RCT runs 24 weeks with a 6 month follow up, and adds a dialect matched condition Wang did not test.",
                   accent=AMBER)
    add_footer(slide, page_num=26)

# ===================== PHASE 5 · REDESIGN (7 slides) =====================

def slide_27_crosswalk(prs):
    slide = slide_blank(prs)
    add_section_label(slide, "PHASE 5 · REDESIGN & RECOMMENDATIONS", color=GREEN)
    add_title(slide, "Every intervention traces back to an HFACS finding.")
    add_subtitle(slide, "Each finding on the left points to one pillar of the redesign on the right.")

    y = Inches(2.85)
    row_h = Inches(0.95)
    gap = Inches(0.1)
    rows = [
        ("L4 · WRONG METRIC · CONVERSATIONAL STYLE UNDERWEIGHTED", "→",  "PERSONA PILLAR (SHIPPED) · Carlos and Elena · Mayer d = 1.3", BLUE),
        ("L3 · BIRDBRAIN OPTIMIZES ITEM CORRECTNESS",              "→",  "DIALOGUE PILLAR (SHIPPED) · scaffolded conversation · Bibauw d = 0.58",   PINK),
        ("L2 · EPISODIC BUFFER NEVER ENGAGED",                     "→",  "CONTEXT PILLAR (SHIPPED) · scenario specific sessions (restaurant, market, meeting)", AMBER),
        ("L1 · WRONG SIGNAL · NO PRODUCTION FEEDBACK",             "→",  "VOICE PILLAR (PROPOSED) · future voice agent · Mayer d = 0.8 · adds a production signal",  RED),
    ]
    for i, (finding, arrow, intervention, color) in enumerate(rows):
        yy = y + i * (row_h + gap)
        add_rect(slide, MARGIN_L, yy, Inches(5.7), row_h, fill_color=CARD, line_color=BORDER, line_width_emu=Emu(6350))
        add_rect(slide, MARGIN_L, yy, Emu(50800), row_h, fill_color=color, line_color=None)
        add_text(slide, Inches(0.9), yy + Inches(0.25), Inches(5.3), row_h - Inches(0.4),
                 finding, size=11, color=TEXT, bold=True, letter_spacing=0.5, spacing=1.35)
        add_text(slide, Inches(6.45), yy + Inches(0.2), Inches(0.6), row_h - Inches(0.3),
                 arrow, size=22, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_rect(slide, Inches(7.15), yy, Inches(5.7), row_h, fill_color=SURFACE, line_color=BORDER, line_width_emu=Emu(6350))
        add_rect(slide, Inches(7.15) + Inches(5.7) - Inches(0.06), yy, Inches(0.06), row_h, fill_color=color, line_color=None)
        add_text(slide, Inches(7.35), yy + Inches(0.25), Inches(5.3), row_h - Inches(0.4),
                 intervention, size=11, color=TEXT_SOFT, spacing=1.35)
    add_footer(slide, page_num=27)

def slide_28_pillars(prs):
    slide = slide_blank(prs)
    add_section_label(slide, "PHASE 5 · REDESIGN", color=GREEN)
    add_title(slide, "Three pillars shipped. Each backed by a meta analysis or a replicated finding.")

    y = Inches(2.8)
    h = Inches(3.7)
    w = Inches(3.95)
    gap = Inches(0.15)
    pillars = [
        ("DIALOGUE WORKS", BLUE,
         "Bibauw et al. (2022)\nd = 0.58 across 17 studies and 100 comparisons.\n\n"
         "A medium meta analytic effect for dialogue based language practice over traditional methods.\n\n"
         "\"System guided interaction outperformed fully free conversation.\" That is the design guidance for our scaffolded agent."),
        ("PERSONA MATTERS", PINK,
         "Mayer (2021)\nConversational style: d = 1.3 across 10 of 10 experiments.\n\n"
         "One of the largest effect sizes in multimedia learning.\n\n"
         "This is not UX decoration. Mayer's research shows conversational style is the mechanism, not the garnish. (Voice and embodiment are future work.)"),
        ("DIALECT IS THE INPUT", AMBER,
         "Geeslin and Long (2014)\n\n"
         "Real conversation requires learning the regional differences in how people actually speak.\n\n"
         "Mexican Spanish and Castilian Spanish differ in sounds, idioms, and formality. The learner needs real regional input, matched to where they are going.\n\n"
         "Carlos and Elena are regionally authentic, not generic \"Spanish.\""),
    ]
    x = MARGIN_L
    for title, color, body in pillars:
        add_accent_card(slide, x, y, w, h, title=title, body=body, accent=color, body_size=10.5)
        x += w + gap
    add_footer(slide, page_num=28)

def slide_29_carlos(prs):
    slide = slide_blank(prs)
    add_section_label(slide, "PHASE 5 · PROTOTYPE", color=GREEN)
    add_title(slide, "Carlos. The Mexican Spanish tutor.")
    add_subtitle(slide, "Built in the repo. Dialect pill, scenario pill, correction callouts, and severity badges.")

    add_placeholder_image(
        slide, MARGIN_L, Inches(2.7), Inches(7.5), Inches(4.3),
        caption_lines=(
            "What should appear here:\n\n"
            "• A chat window mid conversation with Carlos (Mexican Spanish tutor).\n"
            "• Agent reply bubble in italic Spanish with a green \"ES\" badge next to his name.\n"
            "• Mexican flag visible in the dialect pill at the top.\n"
            "• Scenario pill set to \"At a restaurant.\"\n"
            "• At least one visible correction callout (for example, a gentle redirect on word choice).\n"
            "• Severity badge visible on the correction.\n\n"
            "Capture the whole chat container including the header and the context bar."
        ),
        accent=GREEN
    )
    x = Inches(8.3)
    w = Inches(4.6)
    add_rect(slide, x, Inches(2.7), w, Inches(4.3), fill_color=CARD, line_color=BORDER, line_width_emu=Emu(6350))
    add_rect(slide, x, Inches(2.7), Emu(50800), Inches(4.3), fill_color=GREEN, line_color=None)
    add_text(slide, x + Inches(0.25), Inches(2.85), w - Inches(0.4), Inches(0.4),
             "REGIONAL AUTHENTICITY", size=11, color=GREEN, bold=True, letter_spacing=1.2)
    add_text(slide, x + Inches(0.25), Inches(3.3), w - Inches(0.4), Inches(3.5),
             "Carlos uses:\n\n"
             "• Mexican pronunciation (soft /s/, ustedes, no voseo)\n"
             "• Mexican idioms (chido, órale, qué onda, ahorita)\n"
             "• Restaurant vocabulary aligned with Mexican dining (la cuenta, propina, tortillas over pan)\n"
             "• A warm but corrective tutor tone.\n\n"
             "This is Geeslin and Long (2014) put into practice.",
             size=11, color=TEXT_SOFT, spacing=1.5)
    add_footer(slide, page_num=29)

def slide_30_elena(prs):
    slide = slide_blank(prs)
    add_section_label(slide, "PHASE 5 · PROTOTYPE", color=GREEN)
    add_title(slide, "Elena. The Castilian Spanish tutor.")
    add_subtitle(slide, "Same UI, different dialect. One click in the dropdown swaps the linguistic profile.")

    add_placeholder_image(
        slide, MARGIN_L, Inches(2.7), Inches(7.5), Inches(4.3),
        caption_lines=(
            "What should appear here:\n\n"
            "• The same chat UI switched to Elena (Castilian Spanish tutor).\n"
            "• Spanish flag visible in the dialect pill.\n"
            "• Scenario pill set to \"At the market\" or \"Meeting someone new.\"\n"
            "• Translation toggle visible. Ideally OFF (eye closed) so the Spanish only reply shows. This reinforces the cognitive load management annotation on the HTML.\n"
            "• At least one correction callout with a severity badge.\n"
            "• The \"Switched to Castilian Spanish\" system message visible in the chat history."
        ),
        accent=PINK
    )
    x = Inches(8.3)
    w = Inches(4.6)
    add_rect(slide, x, Inches(2.7), w, Inches(4.3), fill_color=CARD, line_color=BORDER, line_width_emu=Emu(6350))
    add_rect(slide, x, Inches(2.7), Emu(50800), Inches(4.3), fill_color=PINK, line_color=None)
    add_text(slide, x + Inches(0.25), Inches(2.85), w - Inches(0.4), Inches(0.4),
             "REGIONAL AUTHENTICITY", size=11, color=PINK, bold=True, letter_spacing=1.2)
    add_text(slide, x + Inches(0.25), Inches(3.3), w - Inches(0.4), Inches(3.5),
             "Elena uses:\n\n"
             "• Castilian pronunciation (distinción; /s/ versus /θ/)\n"
             "• Peninsular grammar (vosotros forms)\n"
             "• Castilian idioms (vale, tío and tía, guay)\n"
             "• Peninsular market vocabulary (metric weights, regional produce, Mercadona norms)\n\n"
             "Same framework. One dropdown switch changes the whole regional identity.",
             size=11, color=TEXT_SOFT, spacing=1.5)
    add_footer(slide, page_num=30)

def slide_31_system_map(prs):
    slide = slide_blank(prs)
    add_section_label(slide, "PHASE 5 · SYSTEM MAP", color=GREEN)
    add_title(slide, "The loop changes shape when we change what it measures.")
    add_subtitle(slide, "BEFORE is today's Duolingo. TARGET is the full redesign, including the proposed voice agent.")

    y = Inches(2.6)
    h = Inches(2.0)
    add_rect(slide, MARGIN_L, y, CONTENT_W, h, fill_color=CARD, line_color=BORDER, line_width_emu=Emu(6350))
    add_rect(slide, MARGIN_L, y, Emu(50800), h, fill_color=RED, line_color=None)
    add_text(slide, Inches(0.95), y + Inches(0.2), Inches(3), Inches(0.35),
             "BEFORE (TODAY)", size=11, color=RED, bold=True, letter_spacing=1.5)
    add_text(slide, Inches(0.95), y + Inches(0.7), Inches(12), Inches(0.5),
             "Learner   →   Item   →   XP",
             size=22, color=TEXT, bold=True)
    add_text(slide, Inches(0.95), y + Inches(1.25), Inches(12), Inches(0.5),
             "A closed loop on the wrong variable. The learner's goal (speaking) is not in the circuit.",
             size=12, color=TEXT_SOFT, italic=True)

    y2 = Inches(4.8)
    add_rect(slide, MARGIN_L, y2, CONTENT_W, h, fill_color=CARD, line_color=BORDER, line_width_emu=Emu(6350))
    add_rect(slide, MARGIN_L, y2, Emu(50800), h, fill_color=GREEN, line_color=None)
    add_text(slide, Inches(0.95), y2 + Inches(0.2), Inches(6), Inches(0.35),
             "TARGET (FULL REDESIGN · INCLUDES FUTURE VOICE AGENT)",
             size=11, color=GREEN, bold=True, letter_spacing=1.5)
    add_text(slide, Inches(0.95), y2 + Inches(0.7), Inches(12), Inches(0.5),
             "Learner   ↔   Dialect Agent   ↔   Scaffolded Dialogue   →   bound context   →   production feedback",
             size=15, color=TEXT, bold=True)
    add_text(slide, Inches(0.95), y2 + Inches(1.25), Inches(12), Inches(0.5),
             "The agent sits in the middle. Context binds. Feedback is about what the learner produced, not what they recognized.",
             size=12, color=TEXT_SOFT, italic=True)
    add_footer(slide, page_num=31)

def slide_32_gaps(prs):
    slide = slide_blank(prs)
    add_section_label(slide, "PHASE 5 · GAP ANALYSIS & TRADE OFFS", color=GREEN)
    add_title(slide, "What we did not build, and the trade offs if we did.")
    add_subtitle(slide, "A gap analysis exposes blind spots. The trade offs keep the proposal honest.")

    y = Inches(2.9)
    h = Inches(3.5)
    w = Inches(6.05)
    gap = Inches(0.25)

    gaps_body = (
        "• Voice input and output is not yet implemented.\n"
        "• The content agent (blogs, songs, news) was descoped from the original proposal.\n"
        "• No A/B test was run. This is a proof of concept only.\n"
        "• No eye tracking data yet. The IX Lab is the next step.\n"
        "• No long term retention data on the prototype.\n"
        "• Two dialects shipped (Mexican and Castilian). Spanish alone has dozens."
    )
    tradeoffs_body = (
        "• Regional authenticity vs. voice model availability. The Castilian voice pool is deeper than the Mexican one (applies once voice is built).\n"
        "• Latency vs. response quality. Fast turn taking competes with deeper model reasoning.\n"
        "• Scope vs. feasibility. A four agent vision against a two week build window.\n"
        "• Correction strength vs. learner anxiety. Too harsh a correction kills willingness to communicate (per Wang 2024).\n"
        "• Regional specificity vs. user reach. A Mexican agent underserves an Argentine learner."
    )
    add_accent_card(slide, MARGIN_L, y, w, h,
                    "GAPS · WHAT IS MISSING", gaps_body, accent=RED, body_size=11)
    add_accent_card(slide, MARGIN_L + w + gap, y, w, h,
                    "TRADE OFFS · WHAT WE WOULD WEIGH", tradeoffs_body, accent=AMBER, body_size=11)
    add_footer(slide, page_num=32)

def slide_33_future(prs):
    slide = slide_blank(prs)
    add_section_label(slide, "PHASE 5 · FUTURE WORK", color=GREEN)
    add_title(slide, "Four agents. One coherent architecture.")
    add_subtitle(slide, "What we shipped is one layer. The voice agent closes the Level 1 signal gap.")

    y = Inches(2.7)
    h = Inches(2.9)
    w = Inches(2.95)
    gap = Inches(0.15)
    agents_list = [
        ("CONTENT AGENT", BLUE, "PROPOSED",
         "Surfaces target language material: blogs, songs, news, speeches.\n\nAddresses the Level 2 hole by giving the learner real context to bind to, even outside of a conversation."),
        ("CONVERSATIONAL AGENT", GREEN, "BUILT",
         "Scaffolded dialogue, scenario lock, and severity graded corrections.\n\nBibauw's d = 0.58 evidence put into practice in a working proof of concept."),
        ("DIALECT AGENTS", PINK, "BUILT",
         "Carlos (Mexican) and Elena (Castilian).\n\nMayer's conversational style research and Geeslin's regional variation research put into practice."),
        ("VOICE AGENT", AMBER, "PROPOSED",
         "Speech to text in. Dialect authentic text to speech out. Fast enough turn taking to feel conversational.\n\nCloses the Level 1 signal loop: the learner speaks, the system hears, and the system can finally measure production."),
    ]
    x = MARGIN_L
    for title, color, status, body in agents_list:
        add_rect(slide, x, y, w, h, fill_color=CARD, line_color=BORDER, line_width_emu=Emu(6350))
        add_rect(slide, x, y, w, Emu(38100), fill_color=color, line_color=None)
        add_text(slide, x + Inches(0.2), y + Inches(0.2), w - Inches(0.4), Inches(0.4),
                 title, size=10.5, color=color, bold=True, letter_spacing=1.2)
        status_color = GREEN if status == "BUILT" else AMBER
        add_text(slide, x + Inches(0.2), y + Inches(0.55), w - Inches(0.4), Inches(0.3),
                 status, size=9.5, color=status_color, bold=True, letter_spacing=1.5)
        add_text(slide, x + Inches(0.2), y + Inches(1.0), w - Inches(0.4), h - Inches(1.1),
                 body, size=10.5, color=TEXT_SOFT, spacing=1.45)
        x += w + gap

    add_pull_quote(slide, MARGIN_L, Inches(5.9), CONTENT_W, Inches(1.1),
                   "The voice layer follows three design principles from this course: recovery (fall back to text if speech fails), redundancy (show the correction even when it is spoken), and reduced cognitive load (voice for production, visual for feedback).",
                   accent=GREEN)
    add_footer(slide, page_num=33)

def slide_34_action_plan(prs):
    slide = slide_blank(prs)
    add_section_label(slide, "PHASE 5 · ACTION PLAN", color=GREEN)
    add_title(slide, "Feasible, and actionable.")
    add_subtitle(slide, "Each step converts one gap analysis item into evidence.")

    y = Inches(2.7)
    row_h = Inches(1.3)
    gap = Inches(0.2)
    rows = [
        ("IMMEDIATE · 0 TO 4 WEEKS", BLUE,
         "An eye tracking study at UT's IX Lab. Sharpens the Level 1 attention claim from conceptual to measurable. Where does the learner actually look: at UI ornaments (XP, streak flame) or at linguistic content?"),
        ("NEAR TERM · 1 TO 6 MONTHS", AMBER,
         "A 24 week randomized controlled trial that extends Wang (2024) in three ways: adds a dialect matched condition, extends measurement to week 24, and adds a 6 month follow up for retention. ACTFL speaking interview scored by a blinded rater. Random assignment stratified by baseline proficiency and intended use context."),
        ("LONG TERM · 6 TO 12 MONTHS", PINK,
         "A longitudinal speaking assessment, which is the metric Duolingo's outcome research has never published. Retention of productive skills at 6 and 12 months, with voice agent condition learners as the primary cell."),
    ]
    for i, (label, color, body) in enumerate(rows):
        yy = y + i * (row_h + gap)
        add_rect(slide, MARGIN_L, yy, CONTENT_W, row_h, fill_color=CARD, line_color=BORDER, line_width_emu=Emu(6350))
        add_rect(slide, MARGIN_L, yy, Emu(50800), row_h, fill_color=color, line_color=None)
        add_text(slide, Inches(0.95), yy + Inches(0.2), Inches(5), Inches(0.35),
                 label, size=11, color=color, bold=True, letter_spacing=1.2)
        add_text(slide, Inches(0.95), yy + Inches(0.6), Inches(12), row_h - Inches(0.7),
                 body, size=11, color=TEXT_SOFT, spacing=1.45)
    add_footer(slide, page_num=34)

# ===================== MAIN =====================

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_title(prs)
    slide_02_agenda(prs)
    slide_03_hook(prs)

    slide_04_system(prs)
    slide_05_stressors(prs)
    slide_06_historical(prs)
    slide_07_tested_failures(prs)
    slide_08_why_matters(prs)
    slide_09_phase1_close(prs)

    slide_10_persona(prs)
    slide_11_recognition_vs_production(prs)
    slide_12_hf_lenses(prs)
    slide_13_ideal_vs_failure(prs)

    slide_14_approach(prs)
    slide_15_build_process(prs)
    slide_16_corpus(prs)
    slide_17_rct(prs)
    slide_18_threats(prs)

    slide_19_hfacs_method(prs)
    slide_20_l4_dau(prs)
    slide_21_l4_personalization(prs)
    slide_22_l3(prs)
    slide_23_l2(prs)
    slide_24_l1(prs)
    slide_25_capstone(prs)
    slide_26_against_grain(prs)

    slide_27_crosswalk(prs)
    slide_28_pillars(prs)
    slide_29_carlos(prs)
    slide_30_elena(prs)
    slide_31_system_map(prs)
    slide_32_gaps(prs)
    slide_33_future(prs)
    slide_34_action_plan(prs)

    out = "/Users/PrestonCusick/Desktop/language-learner-agent/docs/language-learner-agent-research-report.pptx"
    prs.save(out)
    print(f"Saved: {out}  ({len(prs.slides)} slides)")

if __name__ == "__main__":
    main()
